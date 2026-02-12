"""Docling ingestion + AI refinement plugin for UI and MCP runtimes."""

from __future__ import annotations

import os
import tempfile
import base64
import json
import asyncio
import requests
import io
import re
import time
import threading
from datetime import datetime
from functools import lru_cache
from typing import Any, Dict, List, Optional, Tuple

from fastapi import APIRouter, File, HTTPException, UploadFile
from pydantic import BaseModel, Field

from dependencies import get_ai_engine

# Reduce noisy Windows-specific Hugging Face cache warnings in runtime logs.
os.environ.setdefault("HF_HUB_DISABLE_SYMLINKS_WARNING", "1")


router = APIRouter(prefix="/api/docling", tags=["docling", "plugins"])

_LEARNED_ALIASES_CACHE: Optional[Dict[str, Dict[str, Dict[str, int]]]] = None
_MODEL_HEALTH_LOCK = threading.RLock()
_TEMPLATE_PROFILES_LOCK = threading.RLock()
_ALIASES_LOCK = threading.RLock()
_DOCLING_PRIMARY_MODEL = "gemini-2.5-flash-lite"
_DOCLING_FALLBACK_MODEL = "gemini-2.5-flash"
_EXCELLENT_MATCH_THRESHOLD = 99.0
_PROFILE_NON_PERMANENT_CAP = 400
_SMART_ACCEPT_SCORE_THRESHOLD = 85.0
_INCLUDE_EXEMPLAR_VALUES_IN_PROMPTS = False


def _project_root_dir() -> str:
    return os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def _learning_paths() -> Tuple[str, str, str]:
    log_dir = os.path.join(_project_root_dir(), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return (
        os.path.join(log_dir, "docling_learning_raw.jsonl"),
        os.path.join(log_dir, "docling_learned_aliases.json"),
        os.path.join(log_dir, "docling_learning_history.jsonl"),
    )


def _template_profiles_path() -> str:
    log_dir = os.path.join(_project_root_dir(), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return os.path.join(log_dir, "docling_template_profiles.json")


def _template_ranking_history_path() -> str:
    log_dir = os.path.join(_project_root_dir(), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return os.path.join(log_dir, "docling_template_ranking_history.jsonl")


def _template_definitions_path() -> str:
    log_dir = os.path.join(_project_root_dir(), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return os.path.join(log_dir, "docling_template_definitions.json")


def _excellent_match_history_path() -> str:
    log_dir = os.path.join(_project_root_dir(), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return os.path.join(log_dir, "docling_excellent_matches.jsonl")


def _template_upload_history_path() -> str:
    log_dir = os.path.join(_project_root_dir(), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return os.path.join(log_dir, "docling_template_upload_history.jsonl")


def _llm_ranking_history_path() -> str:
    log_dir = os.path.join(_project_root_dir(), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return os.path.join(log_dir, "docling_llm_ranking_history.jsonl")


def _model_health_path() -> str:
    log_dir = os.path.join(_project_root_dir(), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return os.path.join(log_dir, "docling_model_health.json")


def _safe_atomic_json_write(path: str, payload: Any) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    last_err: Optional[Exception] = None
    for attempt in range(1, 8):
        tmp_path: Optional[str] = None
        try:
            with tempfile.NamedTemporaryFile(
                mode="w",
                encoding="utf-8",
                delete=False,
                dir=os.path.dirname(path),
                prefix=os.path.basename(path) + ".",
                suffix=".tmp",
            ) as tmp:
                tmp_path = tmp.name
                json.dump(payload, tmp, ensure_ascii=False, indent=2)
            os.replace(tmp_path, path)
            return
        except (PermissionError, OSError) as exc:
            last_err = exc
            # Windows can briefly lock files; retry with short backoff.
            time.sleep(0.05 * attempt)
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass
    if last_err is not None:
        raise last_err


def _normalize_template_key(name: Optional[str]) -> str:
    t = (name or "").strip().lower().replace("-", "_").replace(" ", "_")
    if t in {"lc", "loc", "letterofcredit"}:
        return "letter_of_credit"
    return t


def _load_model_health() -> Dict[str, Any]:
    path = _model_health_path()
    with _MODEL_HEALTH_LOCK:
        if not os.path.exists(path):
            return {"phase_prefs": {}, "stats": {}, "rankings": {}}
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
            if isinstance(data, dict):
                data.setdefault("phase_prefs", {})
                data.setdefault("stats", {})
                data.setdefault("rankings", {})
                return data
        except Exception:
            pass
    return {"phase_prefs": {}, "stats": {}, "rankings": {}}


def _save_model_health(data: Dict[str, Any]) -> None:
    with _MODEL_HEALTH_LOCK:
        data.setdefault("phase_prefs", {})
        data.setdefault("stats", {})
        data.setdefault("rankings", {})
        _safe_atomic_json_write(_model_health_path(), data)


def _record_model_health(phase: str, model: str, success: bool, error: Optional[str] = None) -> None:
    if not model:
        return
    with _MODEL_HEALTH_LOCK:
        data = _load_model_health()
        phase_stats = data.setdefault("stats", {}).setdefault(phase, {})
        st = phase_stats.setdefault(
            model,
            {"success": 0, "fail": 0, "consecutive_fail": 0, "last_error": None, "last_success": None},
        )
        now = datetime.utcnow().isoformat() + "Z"
        if success:
            st["success"] = int(st.get("success", 0)) + 1
            st["consecutive_fail"] = 0
            st["last_success"] = now
        else:
            st["fail"] = int(st.get("fail", 0)) + 1
            st["consecutive_fail"] = int(st.get("consecutive_fail", 0)) + 1
            st["last_error"] = now
            if error:
                st["last_error_message"] = str(error)[:220]
        _save_model_health(data)


def _model_reliability_score(stats: Dict[str, Any]) -> float:
    s = float(stats.get("success", 0))
    f = float(stats.get("fail", 0))
    consec = float(stats.get("consecutive_fail", 0))
    total = s + f
    success_rate = (s / total) if total > 0 else 0.5
    # Penalize frequent and consecutive failure heavily.
    return round((success_rate * 100.0) + (s * 0.3) - (f * 0.6) - (consec * 8.0), 3)


def _phase_history_bonus(phase: str, history_limit: int = 300) -> Dict[str, float]:
    path = _llm_ranking_history_path()
    if not os.path.exists(path):
        return {}
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            lines = [ln.strip() for ln in f if ln.strip()]
    except Exception:
        return {}

    bonus: Dict[str, float] = {}
    for ln in reversed(lines[-history_limit:]):
        try:
            rec = json.loads(ln)
        except Exception:
            continue
        attempts = rec.get("attempts") or []
        if not isinstance(attempts, list):
            continue
        for attempt in attempts:
            if not isinstance(attempt, dict):
                continue
            model = str(attempt.get("model") or "").strip()
            attempt_phase = str(attempt.get("phase") or "").strip().lower()
            if not model or attempt_phase != phase:
                continue
            success = bool(attempt.get("success"))
            raw_score = float(attempt.get("score") or 0.0)
            err = str(attempt.get("error") or "").lower()
            adj = 0.0
            if success:
                adj += 2.0 + min(2.0, raw_score / 50.0)
            else:
                adj -= 3.5
                if "invalid argument" in err or "error:" in err:
                    adj -= 2.0
            bonus[model] = round(float(bonus.get(model, 0.0)) + adj, 3)
    return bonus


def _rank_models_for_phase(data: Dict[str, Any], phase: str) -> List[Dict[str, Any]]:
    phase_stats = (data.get("stats", {}) or {}).get(phase, {}) or {}
    if not isinstance(phase_stats, dict):
        return []
    hist_bonus = _phase_history_bonus(phase)
    rows: List[Dict[str, Any]] = []
    for model, st in phase_stats.items():
        if not isinstance(st, dict):
            continue
        s = float(st.get("success", 0))
        f = float(st.get("fail", 0))
        total = s + f
        success_rate = round((s / total) * 100.0, 2) if total > 0 else 50.0
        reliability = _model_reliability_score(st)
        bonus = float(hist_bonus.get(model, 0.0))
        score = round(reliability + bonus, 3)
        rows.append(
            {
                "model": model,
                "score": score,
                "reliability_score": reliability,
                "history_bonus": bonus,
                "success_rate": success_rate,
                "attempts": int(total),
                "consecutive_fail": int(st.get("consecutive_fail", 0)),
            }
        )
    rows.sort(
        key=lambda x: (
            float(x.get("score", 0.0)),
            float(x.get("success_rate", 0.0)),
            -int(x.get("consecutive_fail", 0)),
        ),
        reverse=True,
    )
    return rows


def _recompute_phase_preferences(phase: str) -> None:
    with _MODEL_HEALTH_LOCK:
        data = _load_model_health()
        ranked_rows = _rank_models_for_phase(data, phase)
        if not ranked_rows:
            return
        primary = ranked_rows[0]["model"] if ranked_rows else None
        secondary = ranked_rows[1]["model"] if len(ranked_rows) > 1 else None
        data.setdefault("phase_prefs", {})[phase] = {
            "primary": primary,
            "secondary": secondary,
            "updated_at": datetime.utcnow().isoformat() + "Z",
        }
        data.setdefault("rankings", {})[phase] = ranked_rows[:20]
        _save_model_health(data)


async def _get_docling_model_candidates(preferred_model: str, phase: str) -> List[str]:
    preferred_model = (preferred_model or "").strip()
    defaults = [preferred_model, _DOCLING_PRIMARY_MODEL, _DOCLING_FALLBACK_MODEL, "openai/gpt-4.1-nano", "mistralai/ministral-3b-2512"]
    engine = get_ai_engine()
    discovered: List[str] = []
    get_names = getattr(engine, "_get_all_model_names", None)
    if callable(get_names):
        try:
            names = await get_names()
            if isinstance(names, list):
                for n in names:
                    if isinstance(n, str) and n.strip():
                        discovered.append(n.strip())
        except Exception:
            pass

    # Keep preferences fresh so primary/fallback reflect latest performance.
    try:
        _recompute_phase_preferences(phase)
    except Exception:
        pass
    data = _load_model_health()
    prefs = ((data.get("phase_prefs", {}) or {}).get(phase, {}) or {})
    preferred_primary = prefs.get("primary")
    preferred_secondary = prefs.get("secondary")
    ranking_rows = ((data.get("rankings", {}) or {}).get(phase, []) or [])
    phase_stats = ((data.get("stats", {}) or {}).get(phase, {}) or {})
    ranking_score: Dict[str, float] = {}
    if isinstance(ranking_rows, list):
        for row in ranking_rows:
            if isinstance(row, dict):
                model = str(row.get("model") or "").strip()
                if model:
                    ranking_score[model] = float(row.get("score", 0.0))

    # Build pool with ranked primary/fallback first.
    pool: List[str] = []
    for m in [preferred_primary, preferred_secondary] + defaults + discovered:
        if isinstance(m, str) and m.strip() and m.strip() not in pool:
            pool.append(m.strip())

    def sort_key(model: str) -> Tuple[float, int]:
        if model in ranking_score:
            return (float(ranking_score.get(model, 0.0)), 0)
        st = phase_stats.get(model, {}) if isinstance(phase_stats, dict) else {}
        score = _model_reliability_score(st if isinstance(st, dict) else {})
        # Prefer models with fewer consecutive failures when score ties.
        cf = int((st or {}).get("consecutive_fail", 0)) if isinstance(st, dict) else 0
        return (score, -cf)

    ordered = sorted(pool, key=sort_key, reverse=True)
    # Always honor explicitly requested model first; ranking is used for fallback order.
    if preferred_model:
        ordered = [preferred_model] + [m for m in ordered if m != preferred_model]

    return ordered[:20]


def _normalize_text_tokens(text: str) -> List[str]:
    raw = re.findall(r"[a-z0-9]{3,}", (text or "").lower())
    blocked = {
        "invoice", "letter", "credit", "date", "number", "total", "amount",
        "page", "from", "with", "this", "that", "under", "terms", "goods",
    }
    out: List[str] = []
    for tok in raw:
        if tok in blocked:
            continue
        if tok.isdigit() and len(tok) < 4:
            continue
        out.append(tok)
    return out


def _extract_profile_anchors(template: str, extraction: Dict[str, Any]) -> List[str]:
    anchors: List[str] = []
    def collect_tokens(value: Any) -> None:
        if isinstance(value, str) and value.strip():
            anchors.extend(_normalize_text_tokens(value))
            return
        if isinstance(value, (int, float)):
            anchors.extend(_normalize_text_tokens(str(value)))
            return
        if isinstance(value, list):
            for item in value:
                collect_tokens(item)
            return
        if isinstance(value, dict):
            for _, v in value.items():
                collect_tokens(v)
            return

    # Generic extraction for any template structure.
    collect_tokens(extraction)

    # Template-specific boosts.
    if template == "invoice":
        line_items = extraction.get("line_items") or []
        if isinstance(line_items, list):
            for item in line_items[:10]:
                if isinstance(item, dict):
                    desc = item.get("description")
                    if isinstance(desc, str):
                        anchors.extend(_normalize_text_tokens(desc))
    if template == "letter_of_credit":
        docs = ((extraction.get("terms") or {}).get("documents_required") or [])
        if isinstance(docs, list):
            for d in docs[:10]:
                if isinstance(d, str):
                    anchors.extend(_normalize_text_tokens(d))
    # Unique preserve order.
    seen = set()
    uniq: List[str] = []
    for a in anchors:
        if a in seen:
            continue
        seen.add(a)
        uniq.append(a)
    return uniq[:80]


def _template_prompt_anchors(template: str) -> List[str]:
    anchors: List[str] = []
    defs = _load_template_definitions()
    cfg = defs.get(template) or {}
    schema = cfg.get("schema") or {}
    guidance = str(cfg.get("guidance") or "")
    aliases = cfg.get("aliases") or []

    def collect(value: Any) -> None:
        if isinstance(value, dict):
            for k, v in value.items():
                anchors.extend(_normalize_text_tokens(str(k)))
                collect(v)
            return
        if isinstance(value, list):
            for item in value[:8]:
                collect(item)
            return
        if isinstance(value, str):
            anchors.extend(_normalize_text_tokens(value))

    collect(schema)
    anchors.extend(_normalize_text_tokens(guidance))
    if isinstance(aliases, list):
        for a in aliases:
            anchors.extend(_normalize_text_tokens(str(a)))

    # Add learned profile anchors from best-ranked profiles to preserve key lines.
    profiles = _load_template_profiles()
    ranked = []
    for p in profiles:
        if _normalize_learning_template(p.get("template")) != template:
            continue
        ranked.append(
            (
                float(p.get("best_json_score", p.get("avg_json_score", 0.0)) or 0.0),
                int(p.get("seen_count", 0) or 0),
                p.get("anchors") or [],
            )
        )
    ranked.sort(key=lambda x: (x[0], x[1]), reverse=True)
    for _, _, arr in ranked[:3]:
        if isinstance(arr, list):
            for t in arr[:40]:
                anchors.extend(_normalize_text_tokens(str(t)))

    seen = set()
    uniq: List[str] = []
    for a in anchors:
        if a in seen:
            continue
        seen.add(a)
        uniq.append(a)
    return uniq[:200]


def _compact_markdown_for_template(
    template_name: str,
    markdown: str,
    max_chars: int = 12000,
) -> Tuple[str, Dict[str, Any]]:
    src = (markdown or "").strip()
    if not src:
        return "", {"source_chars": 0, "compacted_chars": 0, "compacted": False}

    if len(src) <= max_chars:
        return src, {"source_chars": len(src), "compacted_chars": len(src), "compacted": False}

    template = _resolve_template_name(template_name)
    anchor_set = set(_template_prompt_anchors(template))

    lines = [ln.rstrip() for ln in src.splitlines()]
    scored: List[Tuple[float, int, str]] = []
    for idx, ln in enumerate(lines):
        txt = (ln or "").strip()
        if not txt:
            continue
        toks = _normalize_text_tokens(txt)
        overlap = float(len([t for t in toks if t in anchor_set]))
        has_num = 1.0 if re.search(r"\d", txt) else 0.0
        is_table = 1.0 if "|" in txt else 0.0
        is_header = 1.0 if txt.startswith("#") or txt.endswith(":") else 0.0
        score = (overlap * 4.0) + (has_num * 1.3) + (is_table * 1.8) + (is_header * 0.9)
        scored.append((score, idx, txt))

    # Always keep a small head/tail window for context continuity.
    keep_indices = set(i for i in range(min(20, len(lines))) if lines[i].strip())
    for i in range(max(0, len(lines) - 8), len(lines)):
        if lines[i].strip():
            keep_indices.add(i)

    # Then add best-scoring lines until budget is met.
    scored.sort(key=lambda x: (x[0], -x[1]), reverse=True)
    current_chars = sum(len(lines[i]) + 1 for i in sorted(keep_indices))
    for _, idx, txt in scored:
        if idx in keep_indices:
            continue
        projected = current_chars + len(txt) + 1
        if projected > max_chars:
            continue
        keep_indices.add(idx)
        current_chars = projected
        if current_chars >= max_chars:
            break

    kept = [lines[i] for i in sorted(keep_indices)]
    compacted = "\n".join(kept).strip()
    if len(compacted) > max_chars:
        compacted = compacted[:max_chars].rstrip()
    return compacted, {
        "source_chars": len(src),
        "compacted_chars": len(compacted),
        "compacted": len(compacted) < len(src),
        "kept_lines": len(kept),
        "total_lines": len(lines),
    }


def _jaccard_similarity(a: List[str], b: List[str]) -> float:
    sa = set(a)
    sb = set(b)
    if not sa or not sb:
        return 0.0
    inter = len(sa.intersection(sb))
    uni = len(sa.union(sb))
    return float(inter) / float(uni) if uni else 0.0


def _load_template_profiles() -> List[Dict[str, Any]]:
    with _TEMPLATE_PROFILES_LOCK:
        path = _template_profiles_path()
        if not os.path.exists(path):
            return []
        try:
            with open(path, "r", encoding="utf-8") as f:
                obj = json.load(f)
            if isinstance(obj, list):
                return obj
            return []
        except Exception:
            return []


def _save_template_profiles(profiles: List[Dict[str, Any]]) -> None:
    with _TEMPLATE_PROFILES_LOCK:
        _safe_atomic_json_write(_template_profiles_path(), profiles)


def _upsert_template_profile(
    template: str,
    extraction: Dict[str, Any],
    source: Optional[str],
    match_score: float,
) -> Optional[str]:
    if not isinstance(extraction, dict):
        return None
    anchors = _extract_profile_anchors(template, extraction)
    if not anchors:
        return None
    profiles = _load_template_profiles()
    best_idx = -1
    best_sim = 0.0
    for i, p in enumerate(profiles):
        if p.get("template") != template:
            continue
        sim = _jaccard_similarity(anchors, p.get("anchors") or [])
        if sim > best_sim:
            best_sim = sim
            best_idx = i

    now = datetime.utcnow().isoformat() + "Z"
    if best_idx >= 0 and best_sim >= 0.60:
        prof = profiles[best_idx]
        prof["seen_count"] = int(prof.get("seen_count", 1)) + 1
        old_avg = float(prof.get("avg_match_score", 0.0))
        cnt = int(prof["seen_count"])
        quality = round(float(match_score), 2)
        prof["avg_match_score"] = round(((old_avg * (cnt - 1)) + float(match_score)) / cnt, 2)
        old_avg_json = float(prof.get("avg_json_score", prof.get("avg_match_score", 0.0)))
        prof["avg_json_score"] = round(((old_avg_json * (cnt - 1)) + quality) / cnt, 2)
        best_json = float(prof.get("best_json_score", 0.0))
        if quality >= best_json:
            prof["best_json_score"] = quality
            prof["best_extraction"] = extraction
            prof["best_extraction_updated_at"] = now
        if quality >= _EXCELLENT_MATCH_THRESHOLD:
            prof["permanent"] = True
        prof["last_json_score"] = quality
        # Merge anchors (keep most recent first).
        merged = anchors + [a for a in (prof.get("anchors") or []) if a not in anchors]
        prof["anchors"] = merged[:80]
        prof["updated_at"] = now
        if source:
            prof["sample_source"] = source
        profile_id = str(prof.get("id"))
    else:
        profile_id = f"{template}-{int(datetime.utcnow().timestamp()*1000)}"
        profiles.append(
            {
                "id": profile_id,
                "template": template,
                "created_at": now,
                "updated_at": now,
                "seen_count": 1,
                "avg_match_score": round(float(match_score), 2),
                "avg_json_score": round(float(match_score), 2),
                "best_json_score": round(float(match_score), 2),
                "last_json_score": round(float(match_score), 2),
                "best_extraction": extraction,
                "best_extraction_updated_at": now,
                "permanent": bool(float(match_score) >= _EXCELLENT_MATCH_THRESHOLD),
                "anchors": anchors,
                "sample_source": source,
            }
        )

    # Keep near-permanent excellent profiles; rotate only non-permanent pool.
    permanent = [p for p in profiles if bool(p.get("permanent"))]
    non_permanent = [p for p in profiles if not bool(p.get("permanent"))]
    non_permanent = sorted(non_permanent, key=lambda x: x.get("updated_at", ""), reverse=True)[:_PROFILE_NON_PERMANENT_CAP]
    profiles = permanent + non_permanent
    _save_template_profiles(profiles)
    return profile_id


def _rank_template_profiles(template: str, markdown: str, top_k: int = 10) -> List[Dict[str, Any]]:
    profiles = _load_template_profiles()
    md_tokens = _normalize_text_tokens(markdown or "")
    if not md_tokens:
        return []
    ranked: List[Dict[str, Any]] = []
    for p in profiles:
        if p.get("template") != template:
            continue
        anchors = p.get("anchors") or []
        sim = _jaccard_similarity(md_tokens, anchors)
        if sim <= 0:
            continue
        score = round(min(100.0, sim * 100.0), 2)
        best_json_score = float(p.get("best_json_score", p.get("avg_match_score", 0.0)) or 0.0)
        avg_json_score = float(p.get("avg_json_score", p.get("avg_match_score", 0.0)) or 0.0)
        quality_signal = max(best_json_score, avg_json_score)
        rank_score = round((score * 0.70) + (quality_signal * 0.30), 2)
        ranked.append(
            {
                "profile_id": p.get("id"),
                "template": template,
                "match_score": score,
                "rank_score": rank_score,
                "match_tier": _tier_from_score(score),
                "seen_count": int(p.get("seen_count", 0)),
                "permanent": bool(p.get("permanent")),
                "avg_match_score": float(p.get("avg_match_score", 0.0)),
                "avg_json_score": avg_json_score,
                "best_json_score": best_json_score,
                "sample_source": p.get("sample_source"),
            }
        )
    ranked.sort(
        key=lambda x: (
            -int(bool(x.get("permanent", False))),
            -float(x.get("rank_score", 0.0)),
            -float(x.get("match_score", 0.0)),
            -float(x.get("best_json_score", 0.0)),
            -int(x.get("seen_count", 0)),
        )
    )
    return ranked[: max(1, min(int(top_k), 50))]


def _best_template_profile_extraction(template: str) -> Optional[Dict[str, Any]]:
    target = _normalize_learning_template(template)
    profiles = _load_template_profiles()
    candidates: List[Dict[str, Any]] = []
    for p in profiles:
        t = _normalize_learning_template(p.get("template"))
        if t != target:
            continue
        best = p.get("best_extraction")
        if not isinstance(best, dict) or not best:
            continue
        candidates.append(p)
    if not candidates:
        return None
    candidates.sort(
        key=lambda p: (
            int(bool(p.get("permanent"))),
            float(p.get("best_json_score") or p.get("avg_json_score") or p.get("avg_match_score") or 0.0),
            str(p.get("updated_at") or ""),
        ),
        reverse=True,
    )
    top = candidates[0]
    extraction = json.loads(json.dumps(top.get("best_extraction")))
    notes = extraction.get("extraction_notes") if isinstance(extraction.get("extraction_notes"), list) else []
    notes = list(notes) + ["template_profile_fallback"]
    extraction["extraction_notes"] = notes
    extraction["fallback_profile_id"] = top.get("id")
    extraction["fallback_template"] = target
    return {"profile_id": top.get("id"), "extraction": extraction}


def _append_template_ranking_history_record(
    template: str,
    parser: str,
    source: Optional[str],
    rankings: List[Dict[str, Any]],
    best_profile: Optional[Dict[str, Any]],
) -> None:
    path = _template_ranking_history_path()
    rec = {
        "ts": datetime.utcnow().isoformat() + "Z",
        "template": template,
        "parser": parser,
        "source": source,
        "best_profile": best_profile,
        "rankings": rankings[:10],
    }
    with open(path, "a", encoding="utf-8") as f:
        f.write(json.dumps(rec, ensure_ascii=False) + "\n")


def _append_excellent_match_record(
    template: str,
    source: Optional[str],
    profile_id: Optional[str],
    match_score: float,
    extraction: Optional[Dict[str, Any]],
) -> None:
    if float(match_score) < _EXCELLENT_MATCH_THRESHOLD:
        return
    path = _excellent_match_history_path()
    rec = {
        "ts": datetime.utcnow().isoformat() + "Z",
        "template": _normalize_learning_template(template),
        "source": source,
        "profile_id": profile_id,
        "match_score": round(float(match_score), 2),
        "extraction": extraction if isinstance(extraction, dict) else None,
    }
    with open(path, "a", encoding="utf-8") as f:
        f.write(json.dumps(rec, ensure_ascii=False) + "\n")


def _append_template_upload_history(
    action: str,
    template: str,
    source: Optional[str] = None,
    details: Optional[Dict[str, Any]] = None,
) -> None:
    path = _template_upload_history_path()
    rec = {
        "ts": datetime.utcnow().isoformat() + "Z",
        "action": (action or "unknown").strip().lower(),
        "template": _normalize_learning_template(template),
        "source": source,
        "details": details or {},
    }
    with open(path, "a", encoding="utf-8") as f:
        f.write(json.dumps(rec, ensure_ascii=False) + "\n")


def _load_learned_aliases() -> Dict[str, Dict[str, Dict[str, int]]]:
    global _LEARNED_ALIASES_CACHE
    if _LEARNED_ALIASES_CACHE is not None:
        return _LEARNED_ALIASES_CACHE

    _, aliases_path, _ = _learning_paths()
    if not os.path.exists(aliases_path):
        _LEARNED_ALIASES_CACHE = {}
        return _LEARNED_ALIASES_CACHE

    try:
        with open(aliases_path, "r", encoding="utf-8") as f:
            loaded = json.load(f)
        if isinstance(loaded, dict):
            _LEARNED_ALIASES_CACHE = loaded
        else:
            _LEARNED_ALIASES_CACHE = {}
    except Exception:
        _LEARNED_ALIASES_CACHE = {}
    return _LEARNED_ALIASES_CACHE


def _save_learned_aliases(payload: Dict[str, Dict[str, Dict[str, int]]]) -> None:
    with _ALIASES_LOCK:
        _, aliases_path, _ = _learning_paths()
        _safe_atomic_json_write(aliases_path, payload)


def _normalize_alias(raw: str) -> str:
    cleaned = re.sub(r"[^a-z0-9\s/#&()._-]", " ", (raw or "").lower())
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" -:_")
    return cleaned[:80].strip()


def _is_date_like(value: Optional[str]) -> bool:
    if not value:
        return False
    s = value.strip()
    return bool(
        re.search(r"\b[0-3]?\d[\/\-.][01]?\d[\/\-.](?:\d{2}|\d{4})\b", s)
        or re.search(r"\b(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\s+[0-3]?\d[,]?\s+\d{4}\b", s, flags=re.IGNORECASE)
    )


def _first_date_in_text(value: Optional[str]) -> Optional[str]:
    if not value:
        return None
    m = re.search(r"\b([0-3]?\d[\/\-.][01]?\d[\/\-.](?:\d{2}|\d{4}))\b", value)
    if m:
        return m.group(1)
    m2 = re.search(r"\b((?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\s+[0-3]?\d[,]?\s+\d{4})\b", value, flags=re.IGNORECASE)
    return m2.group(1) if m2 else None


def _is_plausible_invoice_number(value: Optional[str]) -> bool:
    if not value:
        return False
    s = value.strip()
    if len(s) < 5:
        return False
    if re.fullmatch(r"\d{1,4}", s):
        return False
    digits = re.sub(r"\D", "", s)
    return len(digits) >= 4


def _alias_regex(alias: str) -> str:
    return re.escape(alias).replace(r"\ ", r"\s+")


def _learned_labels(template: str, field: str, min_count: int = 1, max_labels: int = 10) -> List[str]:
    store = _load_learned_aliases()
    field_map = ((store.get(template) or {}).get(field) or {})
    labels: List[Tuple[str, int]] = []
    for k, v in field_map.items():
        kk = _normalize_alias(k)
        if not kk:
            continue
        if re.search(r"\b(no\.?\s+of\s+bags?|goods?\s+consigned|number\s+and\s+kind\s+of\s+packages?|for\s+official\s+use)\b", kk):
            continue
        if re.match(r"^\d+[\.\)]\s*", kk):
            continue
        if int(v) >= min_count:
            labels.append((kk, int(v)))
    labels.sort(key=lambda x: (-x[1], x[0]))
    return [k for k, _ in labels[:max_labels]]


def _learned_patterns(
    template: str,
    field: str,
    value_capture_regex: str,
) -> List[str]:
    labels = _learned_labels(template, field, min_count=1, max_labels=12)
    if not labels:
        return []
    joined = "|".join(_alias_regex(lbl) for lbl in labels)
    return [rf"(?i)\b(?:{joined})\s*[:\-]?\s*({value_capture_regex})"]


def _update_learned_alias(template: str, field: str, alias: str) -> None:
    normalized = _normalize_alias(alias)
    if not normalized or len(normalized) < 2:
        return
    if normalized in {"invoice", "letter of credit", "lc", "loc", "details"}:
        return
    if re.search(r"\b(no\.?\s+of\s+bags?|goods?\s+consigned|number\s+and\s+kind\s+of\s+packages?|for\s+official\s+use)\b", normalized):
        return
    if re.match(r"^\d+[\.\)]\s*", normalized):
        return

    store = _load_learned_aliases()
    t_map = store.setdefault(template, {})
    f_map = t_map.setdefault(field, {})
    f_map[normalized] = int(f_map.get(normalized, 0)) + 1
    _save_learned_aliases(store)


def _extract_label_from_line(line: str, value_as_text: str) -> Optional[str]:
    if not line or not value_as_text:
        return None
    m = re.match(r"^\s*([A-Za-z][A-Za-z0-9 /#&()._-]{1,80})\s*[:\-]\s*(.+?)\s*$", line.strip())
    if not m:
        return None
    lhs = (m.group(1) or "").strip()
    rhs = (m.group(2) or "").strip()
    if value_as_text.lower() in rhs.lower():
        return lhs
    return None


def _scalar_text(value: Any) -> Optional[str]:
    if value is None:
        return None
    if isinstance(value, bool):
        return "allowed" if value else "not allowed"
    if isinstance(value, (int, float)):
        txt = f"{value}"
        if "." in txt:
            txt = txt.rstrip("0").rstrip(".")
        return txt
    if isinstance(value, str):
        return value.strip() or None
    return None


def _field_value_is_reliable(template: str, field: str, value: Any) -> bool:
    txt = _scalar_text(value)
    if not txt:
        return False
    currency_codes = {"USD", "EUR", "GBP", "INR", "AED", "JPY", "CNY", "SGD", "AUD", "CAD"}

    if template == "invoice":
        if field == "invoice_number":
            return _is_plausible_invoice_number(txt)
        if field in {"invoice_date", "due_date"}:
            return _is_date_like(txt)
        if field in {"subtotal", "discount_total", "tax_total", "shipping_total", "grand_total", "amount_due"}:
            return _to_float_or_none(txt) is not None
        if field == "currency":
            return txt.strip().upper() in currency_codes
        if field == "po_number":
            s = txt.strip()
            return len(s) >= 4 and bool(re.search(r"\d", s))

    if template == "letter_of_credit":
        if field == "lc_number":
            return len(txt.strip()) >= 6
        if field in {"issue_date", "expiry_date", "latest_shipment_date"}:
            return _is_date_like(txt)
        if field == "amount":
            return _to_float_or_none(txt) is not None
        if field == "currency":
            return txt.strip().upper() in currency_codes
        if field in {"partial_shipment_allowed", "transshipment_allowed"}:
            return isinstance(value, bool)

    return len(txt.strip()) >= 2


def _append_learning_raw_sample(
    template: str,
    parser: str,
    smart_model_used: Optional[str],
    markdown: str,
    structured_extraction: Dict[str, Any],
    json_output: Optional[Dict[str, Any]] = None,
) -> None:
    raw_path, _, _ = _learning_paths()
    sample = {
        "ts": datetime.utcnow().isoformat() + "Z",
        "template": template,
        "parser": parser,
        "smart_model_used": smart_model_used,
        "markdown": markdown,
        "json_output_meta": {
            "source": (json_output or {}).get("source"),
            "parser": (json_output or {}).get("parser"),
            "char_count": (json_output or {}).get("char_count"),
            "line_count": (json_output or {}).get("line_count"),
        },
        "structured_extraction": structured_extraction,
    }
    with open(raw_path, "a", encoding="utf-8") as f:
        f.write(json.dumps(sample, ensure_ascii=False) + "\n")


def _append_learning_history_record(
    template: str,
    parser: Optional[str],
    smart_model_used: Optional[str],
    json_output: Optional[Dict[str, Any]],
    extraction: Optional[Dict[str, Any]],
    status: str,
    error: Optional[str] = None,
) -> None:
    _, _, history_path = _learning_paths()
    normalized_template = _normalize_learning_template(template)
    extraction_obj = extraction if isinstance(extraction, dict) else {}
    score = _score_learning_record(normalized_template, extraction_obj) if extraction_obj else 0.0
    record = {
        "ts": datetime.utcnow().isoformat() + "Z",
        "template": normalized_template,
        "status": status,
        "error": (error or "")[:600] if error else None,
        "match_score": score,
        "match_tier": _tier_from_score(score),
        "parser": parser,
        "smart_model_used": smart_model_used,
        "json_output_meta": {
            "source": (json_output or {}).get("source"),
            "parser": (json_output or {}).get("parser"),
            "char_count": (json_output or {}).get("char_count"),
            "line_count": (json_output or {}).get("line_count"),
        },
    }
    with open(history_path, "a", encoding="utf-8") as f:
        f.write(json.dumps(record, ensure_ascii=False) + "\n")


def _learning_candidate_lines(markdown: str, json_output: Optional[Dict[str, Any]] = None) -> List[str]:
    lines: List[str] = []
    for ln in (markdown or "").splitlines():
        s = ln.strip()
        if not s or s.startswith("## Page") or s.startswith("```"):
            continue
        if ":" not in s:
            continue
        if len(s) > 180:
            continue
        lines.append(s)

    preview_lines = ((json_output or {}).get("preview_lines") or []) if isinstance(json_output, dict) else []
    for ln in preview_lines:
        s = str(ln).strip()
        if not s or ":" not in s or len(s) > 180:
            continue
        lines.append(s)

    # De-duplicate preserving order.
    seen = set()
    out: List[str] = []
    for s in lines:
        if s in seen:
            continue
        seen.add(s)
        out.append(s)
    return out


def _learning_quality_ok(template: str, extraction: Dict[str, Any], smart_model_used: Optional[str]) -> bool:
    # Block noisy self-learning from weak fallback outputs.
    model = (smart_model_used or "").strip().lower()
    is_regex_fallback = model.startswith("regex_")

    if template == "invoice":
        header = extraction.get("header") or {}
        inv_no = header.get("invoice_number")
        inv_dt = header.get("invoice_date")
        strong = _is_plausible_invoice_number(inv_no) and _is_date_like(inv_dt)
        return strong if is_regex_fallback else (strong or _is_plausible_invoice_number(inv_no) or _is_date_like(inv_dt))

    if template == "letter_of_credit":
        header = extraction.get("header") or {}
        lc_no = header.get("lc_number")
        issue_dt = header.get("issue_date")
        strong = bool(lc_no and len(str(lc_no).strip()) >= 6) and _is_date_like(issue_dt)
        return strong if is_regex_fallback else (strong or bool(lc_no))

    score = _score_learning_record(template, extraction if isinstance(extraction, dict) else {})
    return score >= (95.0 if is_regex_fallback else 85.0)


def _learn_from_structured_extraction(
    template: str,
    markdown: str,
    extraction: Dict[str, Any],
    smart_model_used: Optional[str] = None,
    json_output: Optional[Dict[str, Any]] = None,
) -> None:
    if not markdown or not isinstance(extraction, dict):
        return
    if not _learning_quality_ok(template, extraction, smart_model_used):
        return
    lines = _learning_candidate_lines(markdown, json_output=json_output)

    # Learn aliases primarily from header/party/terms scalar fields.
    candidates: List[Tuple[str, Any]] = []
    for section in ("header", "parties", "terms", "financials", "totals"):
        section_obj = extraction.get(section)
        if isinstance(section_obj, dict):
            for field, value in section_obj.items():
                if field == "documents_required":
                    continue
                candidates.append((field, value))

    for field, value in candidates:
        value_text = _scalar_text(value)
        if not value_text:
            continue
        if not _field_value_is_reliable(template, field, value):
            continue
        for ln in lines:
            label = _extract_label_from_line(ln, value_text)
            if label:
                _update_learned_alias(template, field, label)
                break


def _normalize_learning_template(template: Optional[str]) -> str:
    t = _normalize_template_key(template)
    if not t or t in {"all", "*"}:
        return "all"
    return _resolve_template_name(t)


def _reset_learning_state(template: str = "all", clear_raw: bool = True) -> Dict[str, Any]:
    global _LEARNED_ALIASES_CACHE
    normalized_template = _normalize_learning_template(template)
    raw_path, aliases_path, history_path = _learning_paths()
    template_rank_path = _template_ranking_history_path()
    llm_rank_path = _llm_ranking_history_path()
    template_profiles_path = _template_profiles_path()

    aliases = _load_learned_aliases()
    before_templates = sorted(list(aliases.keys()))

    if normalized_template == "all":
        aliases = {}
    else:
        aliases.pop(normalized_template, None)
    _save_learned_aliases(aliases)
    _LEARNED_ALIASES_CACHE = aliases

    raw_total_before = 0
    raw_total_after = 0
    raw_removed = 0
    history_total_before = 0
    history_total_after = 0
    history_removed = 0
    template_rank_before = 0
    template_rank_after = 0
    template_rank_removed = 0
    llm_rank_before = 0
    llm_rank_after = 0
    llm_rank_removed = 0
    template_profiles_before = 0
    template_profiles_after = 0
    model_health_cleared = False
    if clear_raw and os.path.exists(raw_path):
        with open(raw_path, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()
        raw_total_before = len(lines)
        if normalized_template == "all":
            kept_lines: List[str] = []
        else:
            kept_lines = []
            for ln in lines:
                try:
                    rec = json.loads(ln)
                    if _normalize_learning_template(rec.get("template")) == normalized_template:
                        continue
                except Exception:
                    # Keep malformed lines to avoid accidental data loss.
                    pass
                kept_lines.append(ln)
        raw_total_after = len(kept_lines)
        raw_removed = raw_total_before - raw_total_after
        with open(raw_path, "w", encoding="utf-8") as f:
            for ln in kept_lines:
                f.write(ln)
    if clear_raw and os.path.exists(history_path):
        with open(history_path, "r", encoding="utf-8", errors="ignore") as f:
            hlines = f.readlines()
        history_total_before = len(hlines)
        if normalized_template == "all":
            kept_hlines: List[str] = []
        else:
            kept_hlines = []
            for ln in hlines:
                try:
                    rec = json.loads(ln)
                    if _normalize_learning_template(rec.get("template")) == normalized_template:
                        continue
                except Exception:
                    pass
                kept_hlines.append(ln)
        history_total_after = len(kept_hlines)
        history_removed = history_total_before - history_total_after
        with open(history_path, "w", encoding="utf-8") as f:
            for ln in kept_hlines:
                f.write(ln)
    if clear_raw and os.path.exists(template_rank_path):
        with open(template_rank_path, "r", encoding="utf-8", errors="ignore") as f:
            tlines = f.readlines()
        template_rank_before = len(tlines)
        if normalized_template == "all":
            kept_tlines: List[str] = []
        else:
            kept_tlines = []
            for ln in tlines:
                try:
                    rec = json.loads(ln)
                    if _normalize_learning_template(rec.get("template")) == normalized_template:
                        continue
                except Exception:
                    pass
                kept_tlines.append(ln)
        template_rank_after = len(kept_tlines)
        template_rank_removed = template_rank_before - template_rank_after
        with open(template_rank_path, "w", encoding="utf-8") as f:
            for ln in kept_tlines:
                f.write(ln)
    if clear_raw and os.path.exists(llm_rank_path):
        with open(llm_rank_path, "r", encoding="utf-8", errors="ignore") as f:
            llines = f.readlines()
        llm_rank_before = len(llines)
        if normalized_template == "all":
            kept_llines: List[str] = []
        else:
            kept_llines = []
            for ln in llines:
                try:
                    rec = json.loads(ln)
                    if _normalize_learning_template(rec.get("template")) == normalized_template:
                        continue
                except Exception:
                    pass
                kept_llines.append(ln)
        llm_rank_after = len(kept_llines)
        llm_rank_removed = llm_rank_before - llm_rank_after
        with open(llm_rank_path, "w", encoding="utf-8") as f:
            for ln in kept_llines:
                f.write(ln)
    if clear_raw and os.path.exists(template_profiles_path):
        profiles = _load_template_profiles()
        template_profiles_before = len(profiles)
        if normalized_template == "all":
            kept_profiles: List[Dict[str, Any]] = []
        else:
            kept_profiles = [p for p in profiles if _normalize_learning_template(p.get("template")) != normalized_template]
        template_profiles_after = len(kept_profiles)
        _save_template_profiles(kept_profiles)
    if clear_raw and normalized_template == "all":
        try:
            _save_model_health({"phase_prefs": {}, "stats": {}, "rankings": {}})
            model_health_cleared = True
        except Exception:
            model_health_cleared = False

    return {
        "template": normalized_template,
        "clear_raw": clear_raw,
        "aliases_file": aliases_path,
        "raw_file": raw_path,
        "aliases_templates_before": before_templates,
        "aliases_templates_after": sorted(list(aliases.keys())),
        "raw_records_before": raw_total_before,
        "raw_records_after": raw_total_after,
        "raw_records_removed": raw_removed,
        "history_records_before": history_total_before,
        "history_records_after": history_total_after,
        "history_records_removed": history_removed,
        "template_ranking_before": template_rank_before,
        "template_ranking_after": template_rank_after,
        "template_ranking_removed": template_rank_removed,
        "llm_ranking_before": llm_rank_before,
        "llm_ranking_after": llm_rank_after,
        "llm_ranking_removed": llm_rank_removed,
        "template_profiles_before": template_profiles_before,
        "template_profiles_after": template_profiles_after,
        "model_health_cleared": model_health_cleared,
    }


def _to_float_or_none(value: Any) -> Optional[float]:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        try:
            cleaned = re.sub(r"[^\d.\-]", "", value)
            return float(cleaned) if cleaned else None
        except Exception:
            return None
    return None


def _tier_from_score(score: float) -> str:
    if score >= 99:
        return "99% match"
    if score >= 95:
        return "95% match"
    if score >= 90:
        return "90% match"
    return "<90% match"


def _invoice_match_score(extraction: Dict[str, Any]) -> float:
    header = extraction.get("header") or {}
    totals = extraction.get("totals") or {}
    line_items = extraction.get("line_items") or []
    points = 0.0
    total = 100.0

    if _is_plausible_invoice_number(header.get("invoice_number")):
        points += 20
    if _is_date_like(header.get("invoice_date")):
        points += 15
    if header.get("currency"):
        points += 10
    if _to_float_or_none(totals.get("grand_total")) is not None:
        points += 15
    if _to_float_or_none(totals.get("amount_due")) is not None:
        points += 10
    if isinstance(line_items, list) and len(line_items) > 0:
        points += 20

    # Optional enrichment fields.
    opt = 0
    for v in [
        header.get("due_date"),
        header.get("po_number"),
        header.get("seller_name"),
        header.get("buyer_name"),
        totals.get("subtotal"),
        totals.get("tax_total"),
    ]:
        if v not in {None, ""}:
            opt += 1
    points += min(10, opt * 2)

    return round(max(0.0, min(100.0, points / total * 100.0)), 2)


def _lc_match_score(extraction: Dict[str, Any]) -> float:
    header = extraction.get("header") or {}
    parties = extraction.get("parties") or {}
    terms = extraction.get("terms") or {}
    financials = extraction.get("financials") or {}
    points = 0.0
    total = 100.0

    lc_no = (header.get("lc_number") or "").strip() if isinstance(header.get("lc_number"), str) else header.get("lc_number")
    if lc_no and len(str(lc_no)) >= 6:
        points += 20
    if _is_date_like(header.get("issue_date")):
        points += 12
    if _is_date_like(header.get("expiry_date")):
        points += 12
    if header.get("currency"):
        points += 8
    if _to_float_or_none(header.get("amount")) is not None:
        points += 12
    if parties.get("applicant"):
        points += 8
    if parties.get("beneficiary"):
        points += 8
    if terms.get("port_of_loading") or terms.get("port_of_discharge"):
        points += 8

    if isinstance(terms.get("documents_required"), list) and len(terms.get("documents_required") or []) > 0:
        points += 8

    opt = 0
    for v in [
        terms.get("incoterm"),
        terms.get("latest_shipment_date"),
        financials.get("draft_tenor"),
        financials.get("charges"),
    ]:
        if v not in {None, ""}:
            opt += 1
    points += min(4, opt)

    return round(max(0.0, min(100.0, points / total * 100.0)), 2)


def _score_learning_record(template: str, extraction: Dict[str, Any]) -> float:
    if template == "invoice":
        return _invoice_match_score(extraction)
    if template == "letter_of_credit":
        return _lc_match_score(extraction)
    if isinstance(extraction, dict) and extraction:
        return _schema_completeness_score(template, extraction)
    return 0.0


def _load_learning_history(limit: int = 200, template: str = "all") -> Dict[str, Any]:
    normalized_template = _normalize_learning_template(template)
    raw_path, _, history_path = _learning_paths()
    if not os.path.exists(history_path) and not os.path.exists(raw_path):
        return {
            "records": [],
            "summary": {
                "total": 0,
                "avg_match": 0.0,
                "by_tier": {"99% match": 0, "95% match": 0, "90% match": 0, "<90% match": 0},
                "by_status": {"success": 0, "partial": 0, "failed": 0},
            },
        }

    use_history_file = os.path.exists(history_path)
    source_path = history_path if use_history_file else raw_path
    with open(source_path, "r", encoding="utf-8", errors="ignore") as f:
        lines = [ln.strip() for ln in f if ln.strip()]

    parsed: List[Dict[str, Any]] = []
    for ln in reversed(lines):
        try:
            rec = json.loads(ln)
        except Exception:
            continue
        t = _normalize_learning_template(rec.get("template"))
        if normalized_template != "all" and t != normalized_template:
            continue
        if use_history_file:
            score = float(rec.get("match_score") or 0.0)
            tier = rec.get("match_tier") or _tier_from_score(score)
            status = (rec.get("status") or "partial").strip().lower()
            meta = rec.get("json_output_meta") or {}
            error = rec.get("error")
        else:
            extraction = rec.get("structured_extraction") or {}
            score = _score_learning_record(t, extraction if isinstance(extraction, dict) else {})
            tier = _tier_from_score(score)
            status = "success" if isinstance(extraction, dict) and extraction else "partial"
            meta = rec.get("json_output_meta") or {}
            error = None
        parsed.append(
            {
                "ts": rec.get("ts"),
                "template": t,
                "match_score": score,
                "match_tier": tier,
                "status": status,
                "parser": rec.get("parser"),
                "smart_model_used": rec.get("smart_model_used"),
                "source": meta.get("source"),
                "char_count": meta.get("char_count"),
                "error": error,
            }
        )
        if len(parsed) >= max(1, min(int(limit), 1000)):
            break

    tiers = {"99% match": 0, "95% match": 0, "90% match": 0, "<90% match": 0}
    status_counts = {"success": 0, "partial": 0, "failed": 0}
    for p in parsed:
        tiers[p["match_tier"]] = tiers.get(p["match_tier"], 0) + 1
        st = (p.get("status") or "partial").lower()
        if st not in status_counts:
            st = "partial"
        status_counts[st] = status_counts.get(st, 0) + 1
    avg = round(sum(float(p["match_score"]) for p in parsed) / len(parsed), 2) if parsed else 0.0

    return {
        "records": parsed,
        "summary": {
            "total": len(parsed),
            "avg_match": avg,
            "by_tier": tiers,
            "by_status": status_counts,
        },
    }


@lru_cache(maxsize=4)
def _load_docling_converter(quality: str = "fast"):
    try:
        from docling.document_converter import DocumentConverter, PdfFormatOption  # type: ignore
        from docling.datamodel.base_models import InputFormat  # type: ignore
        from docling.datamodel.pipeline_options import PdfPipelineOptions  # type: ignore
    except Exception as exc:
        raise HTTPException(
            status_code=503,
            detail=(
                "Docling is not available. Install it with: "
                "pip install docling"
            ),
        ) from exc

    quality_mode = (quality or "fast").strip().lower()
    if quality_mode not in {"fast", "accurate"}:
        quality_mode = "fast"

    if quality_mode == "fast":
        # Fast profile: best for digital/text PDFs, avoids OCR/table-heavy stages.
        pdf_opts = PdfPipelineOptions()
        pdf_opts.do_ocr = False
        pdf_opts.do_table_structure = False
        pdf_opts.force_backend_text = True
        pdf_opts.document_timeout = 90.0
        pdf_opts.layout_batch_size = 8
        pdf_opts.table_batch_size = 8
        pdf_opts.ocr_batch_size = 8
        return DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_opts),
            }
        )

    return DocumentConverter()


def _convert_source_to_markdown(source: str, quality: str = "fast") -> str:
    converter = _load_docling_converter(quality=quality)
    result = converter.convert(source)
    return result.document.export_to_markdown()


def _guess_mime_from_path(path: str) -> str:
    ext = os.path.splitext((path or "").lower())[1]
    if ext in {".png"}:
        return "image/png"
    if ext in {".jpg", ".jpeg"}:
        return "image/jpeg"
    if ext in {".webp"}:
        return "image/webp"
    if ext in {".gif"}:
        return "image/gif"
    if ext in {".bmp"}:
        return "image/bmp"
    if ext in {".pdf"}:
        return "application/pdf"
    if ext in {".docx"}:
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if ext in {".pptx"}:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if ext in {".xlsx"}:
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if ext in {".html", ".htm"}:
        return "text/html"
    if ext in {".md", ".txt", ".csv", ".json", ".yaml", ".yml", ".xml", ".log"}:
        return "text/plain"
    return "application/octet-stream"


def _structured_json_from_text(text: str, source: str, parser: str, mime_type: str) -> Dict[str, Any]:
    lines = [ln.rstrip() for ln in (text or "").splitlines()]
    non_empty = [ln for ln in lines if ln.strip()]
    return {
        "source": source,
        "parser": parser,
        "mime_type": mime_type,
        "char_count": len(text or ""),
        "line_count": len(lines),
        "non_empty_line_count": len(non_empty),
        "preview_lines": non_empty[:40],
        "text": text or "",
    }


def _extract_pdf_text(path: str) -> str:
    import pypdfium2 as pdfium  # type: ignore

    pdf = pdfium.PdfDocument(path)
    chunks: List[str] = []
    try:
        for i in range(len(pdf)):
            page = pdf[i]
            textpage = page.get_textpage()
            txt = textpage.get_text_range()
            if txt:
                chunks.append(f"## Page {i+1}\n{txt}")
    finally:
        try:
            pdf.close()
        except Exception:
            pass
    return "\n\n".join(chunks).strip()


def _extract_docx_text(path: str) -> str:
    from docx import Document  # type: ignore

    doc = Document(path)
    parts: List[str] = []
    for p in doc.paragraphs:
        if p.text and p.text.strip():
            parts.append(p.text.strip())
    for t_i, table in enumerate(doc.tables, start=1):
        parts.append(f"\n## Table {t_i}")
        for row in table.rows:
            vals = [c.text.strip() for c in row.cells]
            if any(vals):
                parts.append(" | ".join(vals))
    return "\n".join(parts).strip()


def _extract_pptx_text(path: str) -> str:
    from pptx import Presentation  # type: ignore

    prs = Presentation(path)
    parts: List[str] = []
    for s_i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {s_i}")
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                parts.append(text.strip())
    return "\n".join(parts).strip()


def _extract_xlsx_text(path: str) -> str:
    import openpyxl  # type: ignore

    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    parts: List[str] = []
    try:
        for ws in wb.worksheets:
            parts.append(f"## Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if vals:
                    parts.append(" | ".join(vals))
    finally:
        try:
            wb.close()
        except Exception:
            pass
    return "\n".join(parts).strip()


def _extract_html_text(path: str) -> str:
    from bs4 import BeautifulSoup  # type: ignore

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text("\n", strip=True)


def _extract_plain_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _convert_any_local(path: str, mime_type: str) -> Dict[str, Any]:
    ext = os.path.splitext(path.lower())[1]
    parser = "plain_text"
    text = ""
    source = os.path.basename(path)

    if ext == ".pdf":
        parser = "pdfium"
        text = _extract_pdf_text(path)
    elif ext == ".docx":
        parser = "python-docx"
        text = _extract_docx_text(path)
    elif ext == ".pptx":
        parser = "python-pptx"
        text = _extract_pptx_text(path)
    elif ext == ".xlsx":
        parser = "openpyxl"
        text = _extract_xlsx_text(path)
    elif ext in {".html", ".htm"}:
        parser = "beautifulsoup4"
        text = _extract_html_text(path)
    elif ext in {".json"}:
        parser = "json"
        raw = _extract_plain_text(path)
        try:
            obj = json.loads(raw)
            text = json.dumps(obj, indent=2, ensure_ascii=False)
        except Exception:
            text = raw
    elif ext in {".csv", ".md", ".txt", ".py", ".js", ".ts", ".tsx", ".jsx", ".yaml", ".yml", ".xml", ".log"}:
        parser = "plain_text"
        text = _extract_plain_text(path)
    else:
        # Binary/unknown fallback
        parser = "binary_fallback"
        size = os.path.getsize(path)
        text = (
            f"Binary file detected.\n"
            f"Filename: {source}\n"
            f"MIME: {mime_type}\n"
            f"Size: {size} bytes\n"
            "No direct text extraction available for this file type."
        )

    return {
        "markdown": text,
        "json_output": _structured_json_from_text(text, source=source, parser=parser, mime_type=mime_type),
        "parser": parser,
    }


def _is_error_like(text: str) -> bool:
    s = (text or "").strip().lower()
    return s.startswith("error:") or s.startswith("http_error:")


async def _vl_extract_with_fallback(
    image_mime: str,
    image_b64: str,
    primary_model: str,
    fallback_model: str,
) -> Dict[str, Any]:
    prompt = (
        "Extract all visible text from this image as markdown.\n"
        "Rules:\n"
        "- Preserve line breaks and sections.\n"
        "- Do not translate.\n"
        "- Do not add explanations.\n"
        "- If no text is visible, return exactly: NO_TEXT_FOUND\n"
    )
    primary_out = await _call_openrouter_vl_model(primary_model, prompt, image_mime, image_b64)
    if primary_out.strip() and not _is_error_like(primary_out):
        return {"markdown": primary_out, "model_used": primary_model, "fallback_used": False}

    fallback_out = await _call_openrouter_vl_model(fallback_model, prompt, image_mime, image_b64)
    if fallback_out.strip() and not _is_error_like(fallback_out):
        return {"markdown": fallback_out, "model_used": fallback_model, "fallback_used": True}

    raise HTTPException(
        status_code=502,
        detail=f"VL extraction failed on both models. Primary: {primary_out[:220]} | Fallback: {fallback_out[:220]}",
    )


def _smart_vl_prompt_for_template(template_name: str) -> str:
    template = _resolve_template_name(template_name)
    schema = _schema_for_template_name(template)
    if schema is None:
        raise HTTPException(status_code=400, detail=f"Unsupported smart_template '{template_name}'")
    guidance = _template_guidance(template)
    return (
        "You are a strict document extraction engine.\n"
        "Extract structured data directly from the image and return ONLY one JSON object.\n"
        f"{guidance}\n"
        "Output must match this schema keys exactly:\n"
        f"{json.dumps(schema, ensure_ascii=False, indent=2)}\n"
    )


async def _vl_extract_structured_with_fallback(
    image_mime: str,
    image_b64: str,
    template_name: str,
    primary_model: str,
    fallback_model: str,
) -> Dict[str, Any]:
    template = _resolve_template_name(template_name)
    schema = _schema_for_template_name(template)
    if schema is None:
        raise HTTPException(status_code=400, detail=f"Unsupported smart_template '{template_name}'")
    prompt = _smart_vl_prompt_for_template(template)

    attempts: List[Dict[str, Any]] = []

    async def _attempt(model: str) -> Tuple[Optional[Dict[str, Any]], Optional[str], str]:
        out = await _call_openrouter_vl_model(model, prompt, image_mime, image_b64)
        raw = (out or "").strip()
        if not raw or _is_error_like(raw):
            return None, raw[:240], raw
        try:
            parsed = _extract_json_object(raw)
            normalized = _normalize_to_schema(parsed, schema)
            score = _schema_completeness_score(template, normalized)
            attempts.append({
                "phase": "extract",
                "model": model,
                "success": True,
                "score": score,
                "direct": True,
            })
            return normalized, None, raw
        except Exception as exc:
            err = f"invalid_json: {exc}"
            attempts.append({
                "phase": "extract",
                "model": model,
                "success": False,
                "error": err[:220],
                "direct": True,
            })
            return None, err, raw

    primary_data, primary_err, primary_raw = await _attempt(primary_model)
    if primary_data is not None:
        return {
            "data": primary_data,
            "model_used": primary_model,
            "fallback_used": False,
            "llm_attempts": attempts,
            "raw_text": primary_raw,
        }

    fallback_data, fallback_err, fallback_raw = await _attempt(fallback_model)
    if fallback_data is not None:
        return {
            "data": fallback_data,
            "model_used": fallback_model,
            "fallback_used": True,
            "llm_attempts": attempts,
            "raw_text": fallback_raw,
        }

    raise HTTPException(
        status_code=502,
        detail=(
            "Direct structured VL extraction failed on both models. "
            f"Primary: {(primary_err or '')[:220]} | Fallback: {(fallback_err or '')[:220]}"
        ),
    )


async def _extract_pdf_via_vl(path: str, primary_model: str, fallback_model: str, max_pages: int = 5) -> Dict[str, Any]:
    import pypdfium2 as pdfium  # type: ignore

    pdf = pdfium.PdfDocument(path)
    parts: List[str] = []
    used_models: List[str] = []
    fallback_used_any = False
    pages = min(len(pdf), max_pages)
    try:
        for i in range(pages):
            page = pdf[i]
            pil = page.render(scale=2).to_pil()
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            image_b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            res = await _vl_extract_with_fallback(
                "image/png",
                image_b64,
                primary_model,
                fallback_model,
            )
            txt = (res.get("markdown") or "").strip()
            if txt:
                parts.append(f"## Page {i + 1}\n{txt}")
            used_models.append(res.get("model_used") or "")
            fallback_used_any = fallback_used_any or bool(res.get("fallback_used"))
    finally:
        try:
            pdf.close()
        except Exception:
            pass

    merged = "\n\n".join(parts).strip()
    model_used = next((m for m in used_models if m), primary_model)
    return {
        "markdown": merged,
        "model_used": model_used,
        "fallback_used": fallback_used_any,
    }


def _get_openrouter_key() -> str:
    env_key = os.environ.get("OPENROUTER_API_KEY", "").strip()
    if env_key:
        return env_key
    cfg_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "quotas", "quota_openrouter_free.json")
    try:
        with open(cfg_path, "r", encoding="utf-8") as f:
            cfg = json.load(f)
        key = str(cfg.get("api_key", "")).strip()
        return key
    except Exception:
        return ""


def _call_openrouter_vl_model_sync(model: str, prompt: str, image_mime: str, image_b64: str) -> str:
    api_key = _get_openrouter_key()
    if not api_key:
        return "Error: OPENROUTER_API_KEY is missing."

    payload = {
        "model": model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{image_mime};base64,{image_b64}"}},
                ],
            }
        ],
        "temperature": 0,
        "max_tokens": 1200,
    }
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    try:
        r = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=payload, timeout=90)
    except Exception as exc:
        return f"Error: transport failure: {exc}"
    if r.status_code != 200:
        return f"Error: HTTP {r.status_code}: {r.text[:240]}"
    try:
        data = r.json()
        return str(data["choices"][0]["message"]["content"] or "").strip()
    except Exception as exc:
        return f"Error: invalid response format: {exc}"


async def _call_openrouter_vl_model(model: str, prompt: str, image_mime: str, image_b64: str) -> str:
    return await asyncio.to_thread(_call_openrouter_vl_model_sync, model, prompt, image_mime, image_b64)


def _chunk_text(text: str, max_chars: int = 10000) -> List[str]:
    body = (text or "").strip()
    if not body:
        return []
    if len(body) <= max_chars:
        return [body]

    chunks: List[str] = []
    cursor = 0
    n = len(body)
    while cursor < n:
        end = min(cursor + max_chars, n)
        if end < n:
            split = body.rfind("\n\n", cursor, end)
            if split <= cursor:
                split = body.rfind("\n", cursor, end)
            if split <= cursor:
                split = end
            end = split
        chunk = body[cursor:end].strip()
        if chunk:
            chunks.append(chunk)
        cursor = end if end > cursor else cursor + max_chars
    return chunks


def _smart_invoice_schema() -> Dict[str, Any]:
    return {
        "document_type": "invoice",
        "header": {
            "invoice_number": None,
            "invoice_date": None,
            "due_date": None,
            "po_number": None,
            "currency": None,
            "seller_name": None,
            "seller_tax_id": None,
            "buyer_name": None,
            "buyer_tax_id": None,
        },
        "line_items": [
            {
                "line_no": None,
                "description": None,
                "quantity": None,
                "unit": None,
                "unit_price": None,
                "amount": None,
                "tax_rate_percent": None,
                "tax_amount": None,
                "discount_rate_percent": None,
                "discount_amount": None,
            }
        ],
        "totals": {
            "subtotal": None,
            "discount_total": None,
            "tax_total": None,
            "shipping_total": None,
            "grand_total": None,
            "amount_due": None,
        },
        "extraction_notes": [],
    }


def _parse_amount(value: Optional[str]) -> Optional[float]:
    if value is None:
        return None
    cleaned = re.sub(r"[^\d.\-]", "", str(value))
    if not cleaned or cleaned in {"-", ".", "-."}:
        return None
    try:
        return float(cleaned)
    except Exception:
        return None


def _extract_line_items_from_markdown_tables(markdown: str) -> List[Dict[str, Any]]:
    lines = (markdown or "").splitlines()
    blocks: List[List[str]] = []
    current: List[str] = []
    for ln in lines:
        if "|" in ln:
            current.append(ln.strip())
        elif current:
            blocks.append(current)
            current = []
    if current:
        blocks.append(current)

    items: List[Dict[str, Any]] = []
    line_no_counter = 1
    for block in blocks:
        rows: List[List[str]] = []
        for raw in block:
            cells = [c.strip() for c in raw.strip("|").split("|")]
            if len(cells) < 3:
                continue
            if all(re.fullmatch(r":?-{2,}:?", c or "") for c in cells):
                continue
            rows.append(cells)
        if len(rows) < 2:
            continue

        header_idx = -1
        mapping: Dict[str, int] = {}
        for idx, row in enumerate(rows[:3]):
            low = [c.lower() for c in row]
            for i, c in enumerate(low):
                if any(k in c for k in ["description", "item", "product", "goods"]):
                    mapping["description"] = i
                elif any(k in c for k in ["qty", "quantity"]):
                    mapping["quantity"] = i
                elif any(k in c for k in ["unit price", "rate", "price"]):
                    mapping["unit_price"] = i
                elif c in {"unit", "uom"} or " unit " in f" {c} " or "measure" in c:
                    mapping["unit"] = i
                elif any(k in c for k in ["amount", "total", "line total", "value"]):
                    mapping["amount"] = i
                elif c in {"no", "sr", "sr no", "line", "#"} or "line no" in c:
                    mapping["line_no"] = i
            if "description" in mapping and ("amount" in mapping or "unit_price" in mapping):
                header_idx = idx
                break
        if header_idx == -1:
            continue

        for row in rows[header_idx + 1 :]:
            desc_idx = mapping.get("description")
            if desc_idx is None or desc_idx >= len(row):
                continue
            description = row[desc_idx].strip()
            if not description:
                continue

            amount_idx = mapping.get("amount", -1)
            qty_idx = mapping.get("quantity", -1)
            price_idx = mapping.get("unit_price", -1)
            unit_idx = mapping.get("unit", -1)
            line_idx = mapping.get("line_no", -1)

            amount_raw = row[amount_idx].strip() if amount_idx >= 0 and amount_idx < len(row) else None
            qty_raw = row[qty_idx].strip() if qty_idx >= 0 and qty_idx < len(row) else None
            price_raw = row[price_idx].strip() if price_idx >= 0 and price_idx < len(row) else None
            unit_raw = row[unit_idx].strip() if unit_idx >= 0 and unit_idx < len(row) else None
            line_raw = row[line_idx].strip() if line_idx >= 0 and line_idx < len(row) else ""

            parsed_line_no: Optional[int] = None
            if line_raw:
                m_line = re.search(r"\d+", line_raw)
                if m_line:
                    try:
                        parsed_line_no = int(m_line.group(0))
                    except Exception:
                        parsed_line_no = None

            item = {
                "line_no": parsed_line_no if parsed_line_no is not None else line_no_counter,
                "description": description,
                "quantity": _parse_amount(qty_raw),
                "unit": unit_raw or None,
                "unit_price": _parse_amount(price_raw),
                "amount": _parse_amount(amount_raw),
                "tax_rate_percent": None,
                "tax_amount": None,
                "discount_rate_percent": None,
                "discount_amount": None,
            }
            if item["amount"] is None and item["quantity"] is not None and item["unit_price"] is not None:
                item["amount"] = round(float(item["quantity"]) * float(item["unit_price"]), 2)

            items.append(item)
            line_no_counter += 1
    return items


def _extract_invoice_regex_fallback(markdown: str) -> Dict[str, Any]:
    text = markdown or ""
    date_capture = r"(?:[0-3]?\d[\/\-.][01]?\d[\/\-.](?:\d{2}|\d{4})|(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\s+[0-3]?\d[,]?\s+\d{4})"

    def _first_match(patterns: List[str]) -> Optional[str]:
        for pat in patterns:
            m = re.search(pat, text, flags=re.IGNORECASE | re.MULTILINE)
            if m:
                return (m.group(1) or "").strip()
        return None

    header = {
        "invoice_number": _first_match(
            _learned_patterns("invoice", "invoice_number", r"[A-Z0-9\-\/]+")
            + [
                r"\bnumber\s+and\s+date\s+of\s+invoice\b[\s:]*\n?\s*([A-Z0-9\-\/]{5,})\b",
                r"\b(?:invoice\s*(?:no|number|#|id)\s*[:\-]?\s*)([A-Z0-9\-\/]+)\b",
                r"\b(?:inv\s*#?\s*[:\-]?\s*)([A-Z0-9\-\/]+)\b",
            ]
        ),
        "invoice_date": _first_match(
            _learned_patterns("invoice", "invoice_date", date_capture)
            + [
                rf"\bnumber\s+and\s+date\s+of\s+invoice\b[\s:]*\n?\s*[A-Z0-9\-\/]{{5,}}(?:\s*(?:dt|dated|date)\s*[:\-]?\s*)?({date_capture})\b",
                rf"\b(?:invoice\s*date|date)\s*[:\-]?\s*({date_capture})\b",
            ]
        ),
        "due_date": _first_match(
            _learned_patterns("invoice", "due_date", date_capture)
            + [rf"\b(?:due\s*date)\s*[:\-]?\s*({date_capture})\b"]
        ),
        "po_number": _first_match(
            _learned_patterns("invoice", "po_number", r"[A-Z0-9\-\/]+")
            + [
                r"\bpurchase\s*order\s*(?:no|number|#)?\s*[:\-]?\s*([A-Z0-9\-\/]{4,})\b",
                r"\bp\.?\s*o\.?\s*(?:no|number|#)\s*[:\-]?\s*([A-Z0-9\-\/]{4,})\b",
            ]
        ),
        "currency": _first_match([r"\b(USD|EUR|GBP|INR|AED|JPY|CNY|SGD|AUD|CAD)\b"]),
        "seller_name": _first_match(
            _learned_patterns("invoice", "seller_name", r"[^\n|]+")
            + [r"\b(?:seller|supplier|from)\s*[:\-]\s*([^\n|]+)"]
        ),
        "seller_tax_id": _first_match(
            _learned_patterns("invoice", "seller_tax_id", r"[^\n|]+")
            + [r"\b(?:seller\s*tax\s*id|gstin|vat)\s*[:\-]\s*([^\n|]+)"]
        ),
        "buyer_name": _first_match(
            _learned_patterns("invoice", "buyer_name", r"[^\n|]+")
            + [r"\b(?:buyer|bill\s*to|to)\s*[:\-]\s*([^\n|]+)"]
        ),
        "buyer_tax_id": _first_match(
            _learned_patterns("invoice", "buyer_tax_id", r"[^\n|]+")
            + [r"\b(?:buyer\s*tax\s*id)\s*[:\-]\s*([^\n|]+)"]
        ),
    }

    def _match_amount(field_names: str, field_key: Optional[str] = None) -> Optional[float]:
        label_expr = field_names
        if field_key:
            learned = _learned_labels("invoice", field_key, min_count=1, max_labels=12)
            if learned:
                label_expr = f"(?:{field_names}|{'|'.join(_alias_regex(x) for x in learned)})"
        m = re.search(
            rf"\b(?:{label_expr})\b[^\n\d]{{0,40}}([A-Z]{{3}}\s*)?([0-9][0-9,]*\.?\d{{0,2}})\b",
            text,
            flags=re.IGNORECASE,
        )
        if not m:
            return None
        return _parse_amount(m.group(2))

    line_items = _extract_line_items_from_markdown_tables(text)
    subtotal = _match_amount("subtotal", "subtotal")
    discount_total = _match_amount("discount(?:\s*total)?", "discount_total")
    tax_total = _match_amount("tax(?:\s*total)?|vat(?:\s*total)?|gst(?:\s*total)?", "tax_total")
    shipping_total = _match_amount("shipping(?:\s*total)?|freight", "shipping_total")
    grand_total = _match_amount("grand\s*total|total\s*amount|invoice\s*total|f\.?\s*o\.?\s*b\.?\s*value", "grand_total")
    amount_due = _match_amount("amount\s*due|balance\s*due", "amount_due")

    if subtotal is None and line_items:
        subtotal = round(sum(i.get("amount") or 0.0 for i in line_items), 2)
    if grand_total is None and subtotal is not None:
        adjustment = (tax_total or 0.0) + (shipping_total or 0.0) - (discount_total or 0.0)
        grand_total = round(subtotal + adjustment, 2)
    if amount_due is None:
        amount_due = grand_total

    # Post-validate to reduce false positives from OCR noise.
    if not _is_plausible_invoice_number(header.get("invoice_number")):
        header["invoice_number"] = None

    header["invoice_date"] = _first_date_in_text(header.get("invoice_date"))
    header["due_date"] = _first_date_in_text(header.get("due_date"))

    po = (header.get("po_number") or "").strip()
    if po and re.fullmatch(r"\d{1,3}", po):
        header["po_number"] = None

    return {
        "document_type": "invoice",
        "header": header,
        "line_items": line_items,
        "totals": {
            "subtotal": subtotal,
            "discount_total": discount_total,
            "tax_total": tax_total,
            "shipping_total": shipping_total,
            "grand_total": grand_total,
            "amount_due": amount_due,
        },
        "extraction_notes": ["regex_table_fallback"],
    }


def _has_meaningful_invoice_data(data: Optional[Dict[str, Any]]) -> bool:
    if not isinstance(data, dict):
        return False
    header = data.get("header") or {}
    if any(v not in {None, ""} for v in header.values()):
        return True
    line_items = data.get("line_items") or []
    if isinstance(line_items, list) and len(line_items) > 0:
        return True
    totals = data.get("totals") or {}
    if any(v is not None for v in totals.values()):
        return True
    return False


def _to_bool_like(value: Optional[str]) -> Optional[bool]:
    if value is None:
        return None
    s = value.strip().lower()
    if not s:
        return None
    if any(k in s for k in ["not allowed", "not permitted", "no", "false", "prohibited"]):
        return False
    if any(k in s for k in ["allowed", "yes", "true", "permitted"]):
        return True
    return None


def _extract_lc_regex_fallback(markdown: str) -> Dict[str, Any]:
    text = markdown or ""

    def _first_match(patterns: List[str]) -> Optional[str]:
        for pat in patterns:
            m = re.search(pat, text, flags=re.IGNORECASE | re.MULTILINE)
            if m:
                return (m.group(1) or "").strip()
        return None

    currency = _first_match([r"\b(USD|EUR|GBP|INR|AED|JPY|CNY|SGD|AUD|CAD)\b"])
    learned_amount_labels = _learned_labels("letter_of_credit", "amount", min_count=1, max_labels=12)
    amount_label_expr = "amount|credit\\s*amount|lc\\s*amount"
    if learned_amount_labels:
        amount_label_expr = amount_label_expr + "|" + "|".join(_alias_regex(x) for x in learned_amount_labels)
    amount_match = re.search(
        rf"\b(?:{amount_label_expr})\s*[:\-]?\s*([A-Z]{{3}}\s*)?([0-9][0-9,]*\.?\d{{0,2}})\b",
        text,
        flags=re.IGNORECASE,
    )
    amount = _parse_amount(amount_match.group(2)) if amount_match else None
    if amount_match and not currency:
        maybe_curr = (amount_match.group(1) or "").strip().upper()
        currency = maybe_curr or currency

    partial_raw = _first_match(
        _learned_patterns("letter_of_credit", "partial_shipment_allowed", r"[^\n|]+")
        + [r"\bpartial\s*shipment\s*[:\-]?\s*([^\n|]+)"]
    )
    trans_raw = _first_match(
        _learned_patterns("letter_of_credit", "transshipment_allowed", r"[^\n|]+")
        + [r"\btransshipment\s*[:\-]?\s*([^\n|]+)"]
    )
    docs_required: List[str] = []
    lines = (text or "").splitlines()
    for idx, line in enumerate(lines):
        if re.search(r"\bdocuments?\s+required\b", line, flags=re.IGNORECASE):
            for follow in lines[idx + 1 : idx + 15]:
                s = follow.strip()
                if not s:
                    break
                if re.match(r"^[-*]\s+", s):
                    docs_required.append(re.sub(r"^[-*]\s+", "", s).strip())
                    continue
                if re.match(r"^\d+[\).\s]", s):
                    docs_required.append(re.sub(r"^\d+[\).\s]+", "", s).strip())
                    continue
                if ":" in s:
                    break
                docs_required.append(s)
            break

    return {
        "document_type": "letter_of_credit",
        "header": {
            "lc_number": _first_match(
                _learned_patterns("letter_of_credit", "lc_number", r"[A-Z0-9\-\/]+")
                + [r"\b(?:L\/?C|LC|letter\s+of\s+credit)\s*(?:no|number|#)?\s*[:\-]?\s*([A-Z0-9\-\/]+)\b"]
            ),
            "issue_date": _first_match(
                _learned_patterns("letter_of_credit", "issue_date", r"[^\n|]+")
                + [r"\b(?:issue\s*date|date\s*of\s*issue)\s*[:\-]?\s*([^\n|]+)"]
            ),
            "expiry_date": _first_match(
                _learned_patterns("letter_of_credit", "expiry_date", r"[^\n|]+")
                + [r"\b(?:expiry|expiration)\s*date\s*[:\-]?\s*([^\n|]+)"]
            ),
            "expiry_place": _first_match(
                _learned_patterns("letter_of_credit", "expiry_place", r"[^\n|]+")
                + [r"\bexpiry\s*place\s*[:\-]?\s*([^\n|]+)"]
            ),
            "currency": currency,
            "amount": amount,
        },
        "parties": {
            "applicant": _first_match(
                _learned_patterns("letter_of_credit", "applicant", r"[^\n|]+")
                + [r"\bapplicant\s*[:\-]\s*([^\n|]+)"]
            ),
            "beneficiary": _first_match(
                _learned_patterns("letter_of_credit", "beneficiary", r"[^\n|]+")
                + [r"\bbeneficiary\s*[:\-]\s*([^\n|]+)"]
            ),
            "issuing_bank": _first_match(
                _learned_patterns("letter_of_credit", "issuing_bank", r"[^\n|]+")
                + [r"\bissuing\s*bank\s*[:\-]\s*([^\n|]+)"]
            ),
            "advising_bank": _first_match(
                _learned_patterns("letter_of_credit", "advising_bank", r"[^\n|]+")
                + [r"\badvising\s*bank\s*[:\-]\s*([^\n|]+)"]
            ),
            "confirming_bank": _first_match(
                _learned_patterns("letter_of_credit", "confirming_bank", r"[^\n|]+")
                + [r"\bconfirming\s*bank\s*[:\-]\s*([^\n|]+)"]
            ),
        },
        "terms": {
            "incoterm": _first_match(
                _learned_patterns("letter_of_credit", "incoterm", r"[^\n|]+")
                + [r"\bincoterm(?:s)?\s*[:\-]\s*([^\n|]+)", r"\b(FOB|CIF|CFR|EXW|DAP|DDP|FCA)\b"]
            ),
            "latest_shipment_date": _first_match(
                _learned_patterns("letter_of_credit", "latest_shipment_date", r"[^\n|]+")
                + [r"\blatest\s*shipment\s*date\s*[:\-]?\s*([^\n|]+)"]
            ),
            "port_of_loading": _first_match(
                _learned_patterns("letter_of_credit", "port_of_loading", r"[^\n|]+")
                + [r"\bport\s*of\s*loading\s*[:\-]?\s*([^\n|]+)"]
            ),
            "port_of_discharge": _first_match(
                _learned_patterns("letter_of_credit", "port_of_discharge", r"[^\n|]+")
                + [r"\bport\s*of\s*discharge\s*[:\-]?\s*([^\n|]+)"]
            ),
            "partial_shipment_allowed": _to_bool_like(partial_raw),
            "transshipment_allowed": _to_bool_like(trans_raw),
            "documents_required": docs_required,
        },
        "financials": {
            "tolerance_percent": _first_match(
                _learned_patterns("letter_of_credit", "tolerance_percent", r"[0-9]+(?:\.[0-9]+)?%?")
                + [r"\btolerance\s*[:\-]?\s*([0-9]+(?:\.[0-9]+)?%?)"]
            ),
            "available_by": _first_match(
                _learned_patterns("letter_of_credit", "available_by", r"[^\n|]+")
                + [r"\bavailable\s*by\s*[:\-]?\s*([^\n|]+)"]
            ),
            "draft_tenor": _first_match(
                _learned_patterns("letter_of_credit", "draft_tenor", r"[^\n|]+")
                + [r"\bdraft\s*tenor\s*[:\-]?\s*([^\n|]+)"]
            ),
            "charges": _first_match(
                _learned_patterns("letter_of_credit", "charges", r"[^\n|]+")
                + [r"\bcharges?\s*[:\-]?\s*([^\n|]+)"]
            ),
        },
        "extraction_notes": ["regex_lc_fallback"],
    }


def _has_meaningful_lc_data(data: Optional[Dict[str, Any]]) -> bool:
    if not isinstance(data, dict):
        return False
    header = data.get("header") or {}
    parties = data.get("parties") or {}
    terms = data.get("terms") or {}
    financials = data.get("financials") or {}
    if any(v not in {None, ""} for v in header.values()):
        return True
    if any(v not in {None, ""} for v in parties.values()):
        return True
    if any(v not in {None, ""} for k, v in terms.items() if k != "documents_required"):
        return True
    if isinstance(terms.get("documents_required"), list) and len(terms.get("documents_required") or []) > 0:
        return True
    if any(v not in {None, ""} for v in financials.values()):
        return True
    return False


def _smart_lc_schema() -> Dict[str, Any]:
    return {
        "document_type": "letter_of_credit",
        "header": {
            "lc_number": None,
            "issue_date": None,
            "expiry_date": None,
            "expiry_place": None,
            "currency": None,
            "amount": None,
        },
        "parties": {
            "applicant": None,
            "beneficiary": None,
            "issuing_bank": None,
            "advising_bank": None,
            "confirming_bank": None,
        },
        "terms": {
            "incoterm": None,
            "latest_shipment_date": None,
            "port_of_loading": None,
            "port_of_discharge": None,
            "partial_shipment_allowed": None,
            "transshipment_allowed": None,
            "documents_required": [],
        },
        "financials": {
            "tolerance_percent": None,
            "available_by": None,
            "draft_tenor": None,
            "charges": None,
        },
        "extraction_notes": [],
    }


def _default_template_definitions() -> Dict[str, Dict[str, Any]]:
    return {
        "invoice": {
            "schema": _smart_invoice_schema(),
            "guidance": (
                "Target: invoice extraction for headers + line-item grid + discount/tax/subtotals/totals.\n"
                "Rules:\n"
                "- Extract only facts present in the source text.\n"
                "- If a field is missing, use null (or [] for arrays).\n"
                "- Keep monetary and percentage values as strings exactly as seen.\n"
                "- line_items must include only table/grid rows representing billable items.\n"
                "- Do not include prose outside JSON.\n"
            ),
            "aliases": ["invoice"],
        },
        "letter_of_credit": {
            "schema": _smart_lc_schema(),
            "guidance": (
                "Target: letter of credit extraction.\n"
                "Rules:\n"
                "- Extract only explicit facts from source text.\n"
                "- If a field is missing, use null (or [] for arrays).\n"
                "- Keep amounts/dates/terms as strings exactly as seen.\n"
                "- Do not include prose outside JSON.\n"
            ),
            "aliases": ["letter_of_credit", "lc", "loc"],
        },
    }


def _load_template_definitions() -> Dict[str, Dict[str, Any]]:
    path = _template_definitions_path()
    defaults = _default_template_definitions()
    if not os.path.exists(path):
        _safe_atomic_json_write(path, defaults)
        return defaults
    try:
        with open(path, "r", encoding="utf-8") as f:
            raw = json.load(f)
    except Exception:
        _safe_atomic_json_write(path, defaults)
        return defaults

    merged = defaults
    if isinstance(raw, dict):
        for key, value in raw.items():
            norm_key = _normalize_template_key(key)
            if not norm_key or not isinstance(value, dict):
                continue
            schema = value.get("schema")
            if not isinstance(schema, dict):
                continue
            aliases = value.get("aliases")
            aliases_out: List[str] = []
            if isinstance(aliases, list):
                for a in aliases:
                    na = _normalize_template_key(str(a))
                    if na and na not in aliases_out:
                        aliases_out.append(na)
            if norm_key not in aliases_out:
                aliases_out.insert(0, norm_key)
            merged[norm_key] = {
                "schema": schema,
                "guidance": str(value.get("guidance") or "").strip(),
                "aliases": aliases_out,
            }
    _safe_atomic_json_write(path, merged)
    return merged


def _save_template_definitions(defs: Dict[str, Dict[str, Any]]) -> None:
    cleaned: Dict[str, Dict[str, Any]] = {}
    for key, value in (defs or {}).items():
        norm_key = _normalize_template_key(key)
        if not norm_key or not isinstance(value, dict):
            continue
        schema = value.get("schema")
        if not isinstance(schema, dict):
            continue
        aliases = value.get("aliases")
        aliases_out: List[str] = []
        if isinstance(aliases, list):
            for a in aliases:
                na = _normalize_template_key(str(a))
                if na and na not in aliases_out:
                    aliases_out.append(na)
        if norm_key not in aliases_out:
            aliases_out.insert(0, norm_key)
        cleaned[norm_key] = {
            "schema": schema,
            "guidance": str(value.get("guidance") or "").strip(),
            "aliases": aliases_out,
        }
    _safe_atomic_json_write(_template_definitions_path(), cleaned)


def _resolve_template_name(template_name: Optional[str]) -> str:
    norm = _normalize_template_key(template_name)
    if not norm:
        return "all"
    defs = _load_template_definitions()
    if norm in defs:
        return norm
    for key, cfg in defs.items():
        aliases = cfg.get("aliases") or []
        if isinstance(aliases, list) and norm in [str(x) for x in aliases]:
            return key
    return norm


def _schema_for_template_name(template_name: str) -> Optional[Dict[str, Any]]:
    t = _resolve_template_name(template_name)
    if t in {"", "all"}:
        return None
    defs = _load_template_definitions()
    cfg = defs.get(t) or {}
    schema = cfg.get("schema")
    return schema if isinstance(schema, dict) else None


def _schema_from_example(data: Any) -> Any:
    if isinstance(data, dict):
        return {k: _schema_from_example(v) for k, v in data.items()}
    if isinstance(data, list):
        if not data:
            return []
        return [_schema_from_example(data[0])]
    return None


def _guess_template_name_from_payload(payload: Dict[str, Any]) -> str:
    if not isinstance(payload, dict):
        return "custom_template"
    meta = payload.get("document_metadata") if isinstance(payload.get("document_metadata"), dict) else {}
    doc_type = str((meta or {}).get("type") or "").strip()
    low = doc_type.lower()
    if "safta" in low:
        return "safta"
    if "invoice" in low:
        return "invoice"
    if "letter of credit" in low or low in {"lc", "loc"}:
        return "letter_of_credit"
    normalized = _normalize_template_key(doc_type)
    return normalized or "custom_template"


def _register_template_from_example(
    template_name: str,
    exemplar: Dict[str, Any],
    guidance: str = "",
    aliases: Optional[List[str]] = None,
) -> str:
    name = _normalize_template_key(template_name)
    if not name:
        raise ValueError("invalid template name")
    if not isinstance(exemplar, dict) or not exemplar:
        raise ValueError("exemplar must be a non-empty object")
    defs = _load_template_definitions()
    schema = _schema_from_example(exemplar)
    alias_out: List[str] = [name]
    for raw in aliases or []:
        a = _normalize_template_key(raw)
        if a and a not in alias_out:
            alias_out.append(a)
    if name in defs:
        existing_aliases = defs[name].get("aliases") or []
        if isinstance(existing_aliases, list):
            for a in existing_aliases:
                na = _normalize_template_key(str(a))
                if na and na not in alias_out:
                    alias_out.append(na)
    defs[name] = {
        "schema": schema,
        "guidance": (guidance or defs.get(name, {}).get("guidance") or "").strip(),
        "aliases": alias_out,
    }
    _save_template_definitions(defs)
    return name


def _auto_select_template(markdown: str) -> Optional[str]:
    defs = _load_template_definitions()
    candidates = sorted(list(defs.keys()))
    if not candidates:
        return None
    md_tokens = set(_normalize_text_tokens(markdown or ""))
    best_name: Optional[str] = None
    best_score = -1.0
    for name in candidates:
        profile_rows = _rank_template_profiles(name, markdown, top_k=1)
        profile_score = float((profile_rows[0].get("rank_score") if profile_rows else 0.0) or 0.0)
        aliases = defs.get(name, {}).get("aliases") or [name]
        alias_tokens: List[str] = []
        for a in aliases:
            alias_tokens.extend(_normalize_text_tokens(str(a)))
        alias_overlap = 0.0
        if md_tokens and alias_tokens:
            alias_overlap = float(len(md_tokens.intersection(set(alias_tokens)))) * 10.0
        total = profile_score + alias_overlap
        if total > best_score:
            best_score = total
            best_name = name
    return best_name


def _template_guidance(template_name: str) -> str:
    t = _resolve_template_name(template_name)
    defs = _load_template_definitions()
    cfg = defs.get(t) or {}
    g = str(cfg.get("guidance") or "").strip()
    return g or (
        "Target: structured document extraction.\n"
        "Rules:\n"
        "- Extract only explicit facts from source text.\n"
        "- If a field is missing, use null (or [] for arrays).\n"
        "- Keep values as strings exactly as seen unless schema requires otherwise.\n"
        "- Do not include prose outside JSON.\n"
    )


def _clone_schema_default(schema: Any) -> Any:
    if isinstance(schema, dict):
        return {k: _clone_schema_default(v) for k, v in schema.items()}
    if isinstance(schema, list):
        return []
    return None


def _normalize_to_schema(data: Any, schema: Any) -> Any:
    # Keep payload structurally consistent with schema keys/types.
    if isinstance(schema, dict):
        if not isinstance(data, dict):
            return _clone_schema_default(schema)
        out: Dict[str, Any] = {}
        for key, sub_schema in schema.items():
            if key in data:
                out[key] = _normalize_to_schema(data.get(key), sub_schema)
            else:
                out[key] = _clone_schema_default(sub_schema)
        return out
    if isinstance(schema, list):
        if not isinstance(data, list):
            return []
        item_schema = schema[0] if schema else None
        if item_schema is None:
            return data
        return [_normalize_to_schema(item, item_schema) for item in data]
    # Scalar leaf (often None in template schemas): keep extracted value as-is.
    return data


def _smart_prompt_for_template(template_name: str, markdown: str) -> str:
    template = _resolve_template_name(template_name)
    schema = _schema_for_template_name(template)
    if schema is None:
        raise HTTPException(status_code=400, detail=f"Unsupported smart_template '{template_name}'")
    guidance = _template_guidance(template)

    return (
        "You are a strict information extraction engine.\n"
        f"{guidance}\n"
        "Return exactly one JSON object matching this schema (same keys):\n"
        f"{json.dumps(schema, ensure_ascii=False, indent=2)}\n\n"
        "SOURCE MARKDOWN:\n"
        f"{markdown}\n"
    )


def _extract_json_object(raw: Any) -> Dict[str, Any]:
    if isinstance(raw, dict):
        return raw

    text = str(raw or "").strip()
    if not text:
        raise ValueError("Empty model response")

    if text.startswith("```"):
        text = text.strip("`")
        if text.startswith("json"):
            text = text[4:].strip()

    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise ValueError("No JSON object found in model output")

    return json.loads(text[start : end + 1])


class SmartExtractionError(RuntimeError):
    def __init__(self, message: str, attempts: Optional[List[Dict[str, Any]]] = None):
        super().__init__(message)
        self.attempts = attempts or []


def _schema_leaf_stats(data: Any, schema: Any) -> Tuple[int, int]:
    if isinstance(schema, dict):
        total = 0
        hit = 0
        data_obj = data if isinstance(data, dict) else {}
        for k, sub in schema.items():
            t, h = _schema_leaf_stats(data_obj.get(k), sub)
            total += t
            hit += h
        return total, hit
    if isinstance(schema, list):
        # Count list presence as one leaf for completeness.
        if not schema:
            return 1, 1 if isinstance(data, list) and len(data) > 0 else 0
        if not isinstance(data, list) or len(data) == 0:
            return 1, 0
        # Score first item structure + presence.
        t, h = _schema_leaf_stats(data[0], schema[0])
        return max(1, t), max(1, h)
    # scalar leaf
    present = data not in {None, ""} and (not isinstance(data, list) or len(data) > 0)
    return 1, 1 if present else 0


def _schema_completeness_score(template_name: str, data: Dict[str, Any]) -> float:
    schema = _schema_for_template_name(template_name)
    if schema is None:
        return 0.0
    total, hit = _schema_leaf_stats(data, schema)
    if total <= 0:
        return 0.0
    return round((float(hit) / float(total)) * 100.0, 2)


def _append_llm_ranking_record(
    template: str,
    parser: str,
    source: Optional[str],
    attempts: List[Dict[str, Any]],
    selected_model: Optional[str],
    status: str,
    error: Optional[str] = None,
) -> None:
    path = _llm_ranking_history_path()
    rec = {
        "ts": datetime.utcnow().isoformat() + "Z",
        "template": _normalize_learning_template(template),
        "parser": parser,
        "source": source,
        "status": status,
        "selected_model": selected_model,
        "attempts": attempts[:20],
        "error": (error or "")[:800] if error else None,
    }
    with open(path, "a", encoding="utf-8") as f:
        f.write(json.dumps(rec, ensure_ascii=False) + "\n")


def _best_llm_success_score(attempts: List[Dict[str, Any]]) -> float:
    best = 0.0
    for a in attempts or []:
        if not isinstance(a, dict):
            continue
        if not bool(a.get("success")):
            continue
        phase = str(a.get("phase") or "").lower()
        # Score only extraction/repair outputs for structured JSON quality.
        if phase not in {"extract", "repair"}:
            continue
        try:
            score = float(a.get("score") or 0.0)
        except Exception:
            score = 0.0
        if score > best:
            best = score
    return round(best, 2)


def _load_llm_ranking_history(limit: int = 200, template: str = "all") -> Dict[str, Any]:
    path = _llm_ranking_history_path()
    normalized_template = _normalize_learning_template(template)
    if not os.path.exists(path):
        return {
            "records": [],
            "summary": {"total": 0, "by_status": {"success": 0, "partial": 0, "failed": 0}},
        }
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        lines = [ln.strip() for ln in f if ln.strip()]

    recs: List[Dict[str, Any]] = []
    for ln in reversed(lines):
        try:
            rec = json.loads(ln)
        except Exception:
            continue
        t = _normalize_learning_template(rec.get("template"))
        if normalized_template != "all" and t != normalized_template:
            continue
        attempts = rec.get("attempts") or []
        selected_model = rec.get("selected_model")
        attempts = attempts if isinstance(attempts, list) else []
        best_llm_score = _best_llm_success_score(attempts)
        if (
            (rec.get("status") or "").lower() == "success"
            and selected_model
            and "fallback" not in str(selected_model).lower()
            and best_llm_score > 0
            and best_llm_score < _SMART_ACCEPT_SCORE_THRESHOLD
            and _best_template_profile_extraction(t) is not None
        ):
            # Historical normalization: when low-confidence LLM output should have
            # used profile fallback, present it as template_profile_fallback.
            selected_model = "template_profile_fallback"

        recs.append(
            {
                "ts": rec.get("ts"),
                "template": t,
                "status": rec.get("status"),
                "parser": rec.get("parser"),
                "source": rec.get("source"),
                "selected_model": selected_model,
                "attempts": attempts,
                "error": rec.get("error"),
            }
        )
        if len(recs) >= max(1, min(int(limit), 1000)):
            break

    status_counts = {"success": 0, "partial": 0, "failed": 0}
    for r in recs:
        st = (r.get("status") or "partial").lower()
        if st not in status_counts:
            st = "partial"
        status_counts[st] = status_counts.get(st, 0) + 1
    return {"records": recs, "summary": {"total": len(recs), "by_status": status_counts}}


def _seed_text_quality_score(source_markdown: str, seeded_text: str) -> float:
    src_tokens = _normalize_text_tokens(source_markdown or "")
    out_tokens = _normalize_text_tokens(seeded_text or "")
    if not src_tokens or not out_tokens:
        return 0.0
    overlap = _jaccard_similarity(src_tokens, out_tokens) * 100.0
    src_len = max(1, len(source_markdown or ""))
    out_len = len(seeded_text or "")
    ratio = min(float(out_len) / float(src_len), 1.0)
    length_score = ratio * 100.0
    return round((overlap * 0.7) + (length_score * 0.3), 2)


async def _run_llm_seed_text(
    markdown: str,
    template_name: str,
    model_name: str,
) -> Dict[str, Any]:
    engine = get_ai_engine()
    template_norm = _normalize_learning_template(template_name)
    seed_budget = int(os.environ.get("DOCLING_SEED_PROMPT_MAX_CHARS", "14000"))
    compact_source, compact_meta = _compact_markdown_for_template(template_norm, markdown, max_chars=seed_budget)
    prompt = (
        "You are a document text normalizer.\n"
        "Task:\n"
        "- Rewrite OCR markdown into clean plain text key-value lines.\n"
        "- Preserve all factual values exactly.\n"
        "- Remove decorative noise and duplicate lines.\n"
        "- Keep table rows as compact pipe-separated lines if present.\n"
        "Output plain text only.\n\n"
        f"SOURCE:\n{compact_source}\n"
    )
    model_candidates = await _get_docling_model_candidates((model_name or "").strip(), phase="seed_text")

    attempts: List[Dict[str, Any]] = []
    for candidate_model in model_candidates[:8]:
        # Avoid strict response_format hints for provider/model pairs that reject extra arguments.
        raw = await engine.generate_content(candidate_model, prompt)
        text = str(raw or "").strip()
        if text.lower().startswith("error:"):
            _record_model_health("seed_text", candidate_model, success=False, error=text)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "seed_text",
                    "success": False,
                    "score": 0.0,
                    "error": text[:240],
                }
            )
            continue

        if text.startswith("```"):
            text = text.strip("`").strip()
            if text.lower().startswith("markdown"):
                text = text[8:].strip()
            elif text.lower().startswith("text"):
                text = text[4:].strip()

        if len(text) < 40:
            _record_model_health("seed_text", candidate_model, success=False, error="seed output too short")
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "seed_text",
                    "success": False,
                    "score": 0.0,
                    "error": "seed output too short",
                }
            )
            continue

        base_score = _seed_text_quality_score(markdown, text)
        ranked_profiles = _rank_template_profiles(template_norm, text, top_k=1) if template_norm not in {"", "all"} else []
        profile_score = float((ranked_profiles[0].get("match_score") if ranked_profiles else 0.0) or 0.0)
        final_score = round((base_score * 0.7) + (profile_score * 0.3), 2)
        attempts.append(
            {
                "model": candidate_model,
                "phase": "seed_text",
                "success": True,
                "score": final_score,
                "seed_quality": base_score,
                "template_similarity": profile_score,
                "error": None,
                "text": text,
            }
        )
        _record_model_health("seed_text", candidate_model, success=True)

    ranked = sorted(attempts, key=lambda x: (bool(x.get("success")), float(x.get("score", 0.0))), reverse=True)
    best = next((a for a in ranked if a.get("success")), None)
    if best is None:
        _recompute_phase_preferences("seed_text")
        raise SmartExtractionError("seed text generation failed", attempts=attempts)
    _recompute_phase_preferences("seed_text")
    return {
        "text": best.get("text") or markdown,
        "seed_model_used": best.get("model"),
        "seed_attempts": attempts,
        "seed_ranking": ranked,
        "prompt_compaction": compact_meta,
    }


async def _run_smart_template_extraction(
    markdown: str,
    template_name: str,
    model_name: str,
) -> Dict[str, Any]:
    engine = get_ai_engine()
    extract_budget = int(os.environ.get("DOCLING_EXTRACT_PROMPT_MAX_CHARS", "12000"))
    compact_source, compact_meta = _compact_markdown_for_template(template_name, markdown, max_chars=extract_budget)
    prompt = _smart_prompt_for_template(template_name, compact_source)

    model_candidates = await _get_docling_model_candidates((model_name or "").strip(), phase="smart_extract")

    errors: List[str] = []
    attempts: List[Dict[str, Any]] = []
    for candidate_model in model_candidates[:8]:
        # Attempt 1: strict schema extraction
        # Avoid strict response_format hints for provider/model pairs that reject extra arguments.
        raw = await engine.generate_content(candidate_model, prompt)
        raw_text = str(raw or "").strip()
        if raw_text.lower().startswith("error:"):
            errors.append(f"{candidate_model}: {raw_text[:180]}")
            _record_model_health("smart_extract", candidate_model, success=False, error=raw_text)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "extract",
                    "success": False,
                    "score": 0.0,
                    "error": raw_text[:240],
                }
            )
            continue
        try:
            parsed = _extract_json_object(raw)
            score = _schema_completeness_score(template_name, parsed)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "extract",
                    "success": True,
                    "score": score,
                    "error": None,
                }
            )
            _record_model_health("smart_extract", candidate_model, success=True)
            ranked = sorted(attempts, key=lambda x: (bool(x.get("success")), float(x.get("score", 0.0))), reverse=True)
            _recompute_phase_preferences("smart_extract")
            return {
                "data": parsed,
                "smart_model_used": candidate_model,
                "llm_attempts": attempts,
                "llm_ranking": ranked,
                "prompt_compaction": compact_meta,
            }
        except Exception as exc:
            errors.append(f"{candidate_model}: parse failed: {exc}")
            _record_model_health("smart_extract", candidate_model, success=False, error=f"parse failed: {exc}")
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "extract",
                    "success": False,
                    "score": 0.0,
                    "error": f"parse failed: {str(exc)[:200]}",
                }
            )

        # Attempt 2: JSON repair pass on same model
        repair_prompt = (
            "Fix the following model output into one valid JSON object only.\n"
            "Do not add explanations.\n"
            "Output must be strict JSON.\n\n"
            f"Broken output:\n{raw_text}\n"
        )
        # Keep repair call permissive; parse JSON ourselves from model text.
        repaired = await engine.generate_content(candidate_model, repair_prompt)
        repaired_text = str(repaired or "").strip()
        if repaired_text.lower().startswith("error:"):
            errors.append(f"{candidate_model}: repair failed: {repaired_text[:180]}")
            _record_model_health("smart_extract", candidate_model, success=False, error=repaired_text)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "repair",
                    "success": False,
                    "score": 0.0,
                    "error": repaired_text[:240],
                }
            )
            continue
        try:
            parsed_repair = _extract_json_object(repaired)
            score = _schema_completeness_score(template_name, parsed_repair)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "repair",
                    "success": True,
                    "score": score,
                    "error": None,
                }
            )
            _record_model_health("smart_extract", candidate_model, success=True)
            ranked = sorted(attempts, key=lambda x: (bool(x.get("success")), float(x.get("score", 0.0))), reverse=True)
            _recompute_phase_preferences("smart_extract")
            return {
                "data": parsed_repair,
                "smart_model_used": candidate_model,
                "llm_attempts": attempts,
                "llm_ranking": ranked,
                "prompt_compaction": compact_meta,
            }
        except Exception as exc:
            errors.append(f"{candidate_model}: repair parse failed: {exc}")
            _record_model_health("smart_extract", candidate_model, success=False, error=f"repair parse failed: {exc}")
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "repair",
                    "success": False,
                    "score": 0.0,
                    "error": f"repair parse failed: {str(exc)[:200]}",
                }
            )

    _recompute_phase_preferences("smart_extract")
    raise SmartExtractionError(" ; ".join(errors[:8]) if errors else "unknown smart extraction error", attempts=attempts)


async def _run_smart_template_extraction_from_binary(
    binary_bytes: bytes,
    mime_type: str,
    template_name: str,
    model_name: str,
) -> Dict[str, Any]:
    engine = get_ai_engine()
    prompt = _smart_prompt_for_template(template_name, "Use the attached file content as the only source of truth.")
    inline_payload = [{"mime_type": mime_type, "data": base64.b64encode(binary_bytes).decode("ascii")}]

    model_candidates = await _get_docling_model_candidates((model_name or "").strip(), phase="smart_extract")

    errors: List[str] = []
    attempts: List[Dict[str, Any]] = []
    for candidate_model in model_candidates[:8]:
        raw = await engine.generate_content(candidate_model, prompt, images=inline_payload)
        raw_text = str(raw or "").strip()
        if raw_text.lower().startswith("error:"):
            errors.append(f"{candidate_model}: {raw_text[:180]}")
            _record_model_health("smart_extract", candidate_model, success=False, error=raw_text)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "extract_binary",
                    "success": False,
                    "score": 0.0,
                    "error": raw_text[:240],
                }
            )
            continue
        try:
            parsed = _extract_json_object(raw)
            score = _schema_completeness_score(template_name, parsed)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "extract_binary",
                    "success": True,
                    "score": score,
                    "error": None,
                }
            )
            _record_model_health("smart_extract", candidate_model, success=True)
            ranked = sorted(attempts, key=lambda x: (bool(x.get("success")), float(x.get("score", 0.0))), reverse=True)
            _recompute_phase_preferences("smart_extract")
            return {
                "data": parsed,
                "smart_model_used": candidate_model,
                "llm_attempts": attempts,
                "llm_ranking": ranked,
                "prompt_compaction": {"mode": "pdf_direct_binary"},
                "raw_text": raw_text,
            }
        except Exception as exc:
            errors.append(f"{candidate_model}: parse failed: {exc}")
            _record_model_health("smart_extract", candidate_model, success=False, error=f"parse failed: {exc}")
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "extract_binary",
                    "success": False,
                    "score": 0.0,
                    "error": f"parse failed: {str(exc)[:200]}",
                }
            )

        repair_prompt = (
            "Fix the following model output into one valid JSON object only.\n"
            "Do not add explanations.\n"
            "Output must be strict JSON.\n\n"
            f"Broken output:\n{raw_text}\n"
        )
        repaired = await engine.generate_content(candidate_model, repair_prompt)
        repaired_text = str(repaired or "").strip()
        if repaired_text.lower().startswith("error:"):
            errors.append(f"{candidate_model}: repair failed: {repaired_text[:180]}")
            _record_model_health("smart_extract", candidate_model, success=False, error=repaired_text)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "repair_binary",
                    "success": False,
                    "score": 0.0,
                    "error": repaired_text[:240],
                }
            )
            continue
        try:
            parsed_repair = _extract_json_object(repaired)
            score = _schema_completeness_score(template_name, parsed_repair)
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "repair_binary",
                    "success": True,
                    "score": score,
                    "error": None,
                }
            )
            _record_model_health("smart_extract", candidate_model, success=True)
            ranked = sorted(attempts, key=lambda x: (bool(x.get("success")), float(x.get("score", 0.0))), reverse=True)
            _recompute_phase_preferences("smart_extract")
            return {
                "data": parsed_repair,
                "smart_model_used": candidate_model,
                "llm_attempts": attempts,
                "llm_ranking": ranked,
                "prompt_compaction": {"mode": "pdf_direct_binary"},
                "raw_text": repaired_text,
            }
        except Exception as exc:
            errors.append(f"{candidate_model}: repair parse failed: {exc}")
            _record_model_health("smart_extract", candidate_model, success=False, error=f"repair parse failed: {exc}")
            attempts.append(
                {
                    "model": candidate_model,
                    "phase": "repair_binary",
                    "success": False,
                    "score": 0.0,
                    "error": f"repair parse failed: {str(exc)[:200]}",
                }
            )

    _recompute_phase_preferences("smart_extract")
    raise SmartExtractionError(" ; ".join(errors[:8]) if errors else "unknown smart extraction error", attempts=attempts)


class RefineRequest(BaseModel):
    markdown: str
    instruction: str
    model: Optional[str] = None
    template: Optional[str] = None
    chunk_size_chars: Optional[int] = 10000
    use_superpowers: Optional[bool] = False
    superpower_mode: Optional[str] = "review"


class TemplateDefinitionRequest(BaseModel):
    name: str
    schema_payload: Dict[str, Any] = Field(alias="schema")
    guidance: Optional[str] = ""
    aliases: Optional[List[str]] = None


class TemplateExemplarRequest(BaseModel):
    template: Optional[str] = None
    exemplar: Dict[str, Any]
    guidance: Optional[str] = ""
    aliases: Optional[List[str]] = None
    source: Optional[str] = "manual"


@router.post("/convert")
async def convert_document(
    source: Optional[str] = None,
    file: Optional[UploadFile] = File(default=None),
    quality: str = "fast",
    auto_retry: bool = True,
    extractor: str = "auto",
    vl_model: str = "mistralai/ministral-3b-2512",
    vl_fallback_model: str = "openai/gpt-4.1-nano",
    smart_template: str = "none",
    smart_model: str = _DOCLING_PRIMARY_MODEL,
    smart_required: bool = True,
    smart_single_call: bool = True,
    pdf_direct_llm: bool = True,
) -> Dict[str, Any]:
    if not source and not file:
        raise HTTPException(status_code=400, detail="Provide either 'source' or 'file'.")

    temp_path: Optional[str] = None
    try:
        if file is not None:
            filename = file.filename or "uploaded_document"
            suffix = os.path.splitext(filename)[1]
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                temp_path = tmp.name
                data = await file.read()
                tmp.write(data)
            source_to_convert = temp_path
            mime_type = file.content_type or _guess_mime_from_path(filename)
        else:
            source_to_convert = str(source)
            mime_type = _guess_mime_from_path(source_to_convert)

        extractor_mode = (extractor or "auto").strip().lower()
        smart_name = (smart_template or "none").strip().lower()
        smart_auto = smart_name == "auto"
        resolved_template = _resolve_template_name(smart_name) if smart_name not in {"none", "auto"} else smart_name
        smart_schema = _schema_for_template_name(resolved_template) if smart_name not in {"none", "auto"} else None
        smart_enabled = smart_name != "none" and (smart_auto or smart_schema is not None)
        if smart_name not in {"none", "auto"} and not smart_enabled:
            raise HTTPException(
                status_code=400,
                detail=(
                    f"Unknown smart_template '{smart_name}'. "
                    "Create it via /api/docling/templates/definitions first."
                ),
            )
        if extractor_mode == "auto":
            extractor_mode = "vl" if mime_type.startswith("image/") else "local"
        elif extractor_mode == "local" and smart_name != "none" and mime_type.startswith("image/"):
            # Local parser on images returns no meaningful OCR text, so prefer VL for structured extraction.
            extractor_mode = "vl"

        direct_single_call_used = False
        direct_llm_attempts: List[Dict[str, Any]] = []
        direct_structured_extraction: Optional[Dict[str, Any]] = None
        direct_smart_model_used: Optional[str] = None
        direct_prompt_compaction: Optional[Dict[str, Any]] = None

        # Fast path for all binary documents: one direct smart extraction call with attached file bytes.
        # This bypasses OCR/VL markdown generation and seed stage entirely.
        if (
            pdf_direct_llm
            and smart_enabled
            and smart_single_call
            and not smart_auto
            and smart_schema is not None
        ):
            binary_bytes: Optional[bytes] = None
            if file is not None and isinstance(data, (bytes, bytearray)):
                binary_bytes = bytes(data)
            elif os.path.exists(source_to_convert) and os.path.isfile(source_to_convert):
                try:
                    with open(source_to_convert, "rb") as f:
                        binary_bytes = f.read()
                except Exception:
                    binary_bytes = None
            if binary_bytes:
                try:
                    direct_result = await _run_smart_template_extraction_from_binary(
                        binary_bytes=binary_bytes,
                        mime_type=mime_type,
                        template_name=resolved_template,
                        model_name=(smart_model or _DOCLING_PRIMARY_MODEL).strip(),
                    )
                    direct_single_call_used = True
                    direct_llm_attempts = direct_result.get("llm_attempts") or []
                    direct_structured_extraction = direct_result.get("data")
                    direct_smart_model_used = direct_result.get("smart_model_used")
                    direct_prompt_compaction = direct_result.get("prompt_compaction")
                    # Direct mode should return structured JSON only (no markdown/raw text echo).
                    markdown = ""
                    used_model = None
                    fallback_used = False
                    used_quality = "direct"
                    parser_used = "direct-llm"
                    extractor_mode = "direct-llm"
                    json_output = _structured_json_from_text(
                        json.dumps(direct_structured_extraction or {}, ensure_ascii=False),
                        source=(file.filename if file else source_to_convert),
                        parser=parser_used,
                        mime_type=mime_type,
                    )
                except Exception as exc:
                    # Do not route into OCR/VL fallback when smart direct mode is required.
                    if smart_required:
                        raise HTTPException(
                            status_code=422,
                            detail=f"Direct single-call extraction failed: {str(exc)}",
                        ) from exc
                    direct_single_call_used = False

        if not direct_single_call_used and extractor_mode == "vl":
            if not mime_type.startswith("image/"):
                raise HTTPException(
                    status_code=400,
                    detail="VL extractor currently supports images only. Use extractor=docling for PDFs/docs.",
                )
            if file is not None:
                image_bytes = data
            else:
                with open(source_to_convert, "rb") as f:
                    image_bytes = f.read()
            image_b64 = base64.b64encode(image_bytes).decode("ascii")
            # Fast single-call path: image -> structured JSON directly (skip markdown->template second call).
            if smart_enabled and not smart_auto and smart_schema is not None:
                direct_result = await _vl_extract_structured_with_fallback(
                    mime_type, image_b64, resolved_template, vl_model, vl_fallback_model
                )
                direct_single_call_used = True
                direct_llm_attempts = direct_result.get("llm_attempts") or []
                direct_structured_extraction = direct_result.get("data")
                direct_smart_model_used = direct_result.get("model_used")
                # Direct mode should return structured JSON only (no markdown/raw text echo).
                markdown = ""
                used_model = direct_result["model_used"]
                fallback_used = bool(direct_result.get("fallback_used"))
            else:
                vl_result = await _vl_extract_with_fallback(mime_type, image_b64, vl_model, vl_fallback_model)
                markdown = vl_result["markdown"]
                used_model = vl_result["model_used"]
                fallback_used = vl_result["fallback_used"]
            json_output = _structured_json_from_text(
                json.dumps(direct_structured_extraction or {}, ensure_ascii=False) if direct_single_call_used else markdown,
                source=(file.filename if file else source_to_convert),
                parser="vl-ocr",
                mime_type=mime_type,
            )
            used_quality = "vl"
            parser_used = "vl-ocr"
        elif not direct_single_call_used and extractor_mode == "local":
            local_result = _convert_any_local(source_to_convert, mime_type=mime_type)
            markdown = local_result["markdown"]
            json_output = local_result["json_output"]
            used_quality = "local"
            used_model = None
            fallback_used = False
            parser_used = local_result["parser"]
        elif not direct_single_call_used:
            used_quality = quality
            markdown = _convert_source_to_markdown(source_to_convert, quality=quality)
            # Fast mode can return almost nothing for scanned/image PDFs when OCR is disabled.
            if auto_retry and quality.strip().lower() == "fast" and len((markdown or "").strip()) < 80:
                markdown = _convert_source_to_markdown(source_to_convert, quality="accurate")
                used_quality = "accurate"
            used_model = None
            fallback_used = used_quality != quality
            parser_used = "docling"
            json_output = _structured_json_from_text(
                markdown,
                source=(file.filename if file else source_to_convert),
                parser="docling",
                mime_type=mime_type,
            )

        # Prevent false-positive "success with 0 chars".
        if not direct_single_call_used and not (markdown or "").strip():
            # For PDFs, auto-fallback to VL OCR of rendered pages.
            if extractor_mode in {"auto", "local"} and mime_type == "application/pdf":
                vl_pdf = await _extract_pdf_via_vl(source_to_convert, vl_model, vl_fallback_model, max_pages=5)
                markdown = vl_pdf["markdown"]
                used_quality = "vl"
                used_model = vl_pdf["model_used"]
                fallback_used = vl_pdf["fallback_used"]
                parser_used = "pdf-vl-ocr"
                extractor_mode = "vl"
                json_output = _structured_json_from_text(
                    markdown,
                    source=(file.filename if file else source_to_convert),
                    parser=parser_used,
                    mime_type=mime_type,
                )

        if not direct_single_call_used and not (markdown or "").strip():
            raise HTTPException(
                status_code=422,
                detail="No text could be extracted from this file. Try a clearer scan or a different extractor.",
            )

        extraction_markdown = markdown
        llm_seed_attempts: List[Dict[str, Any]] = []
        llm_seed_ranking: List[Dict[str, Any]] = []
        llm_seed_model_used: Optional[str] = None
        llm_seed_prompt_compaction: Optional[Dict[str, Any]] = None
        seed_stage_skipped = False
        if smart_auto:
            chosen_template = _auto_select_template(markdown)
            if chosen_template:
                resolved_template = chosen_template
                smart_schema = _schema_for_template_name(resolved_template)
            else:
                smart_enabled = False
            if smart_required and not smart_enabled:
                raise HTTPException(
                    status_code=422,
                    detail="Auto template selection failed. Upload template profiles/examples first.",
                )
        if smart_enabled and not direct_single_call_used:
            if smart_single_call:
                # Single-LLM mode for smart extraction: skip seed normalization call.
                seed_stage_skipped = True
                extraction_markdown = markdown
                llm_seed_prompt_compaction = {"mode": "bypass", "reason": "smart_single_call"}
            else:
                try:
                    seed_result = await _run_llm_seed_text(
                        markdown=markdown,
                        template_name=resolved_template,
                        model_name=(smart_model or _DOCLING_PRIMARY_MODEL).strip(),
                    )
                    extraction_markdown = seed_result.get("text") or markdown
                    llm_seed_model_used = seed_result.get("seed_model_used")
                    llm_seed_attempts = seed_result.get("seed_attempts") or []
                    llm_seed_ranking = seed_result.get("seed_ranking") or []
                    llm_seed_prompt_compaction = seed_result.get("prompt_compaction")
                except Exception as exc:
                    if isinstance(exc, SmartExtractionError):
                        llm_seed_attempts = exc.attempts or []
                        llm_seed_ranking = sorted(
                            llm_seed_attempts,
                            key=lambda x: (bool(x.get("success")), float(x.get("score", 0.0))),
                            reverse=True,
                        )

        template_rankings: List[Dict[str, Any]] = []
        best_template_profile: Optional[Dict[str, Any]] = None
        if smart_enabled:
            rank_template = resolved_template
            try:
                template_rankings = _rank_template_profiles(rank_template, extraction_markdown, top_k=10)
                best_template_profile = template_rankings[0] if template_rankings else None
                _append_template_ranking_history_record(
                    template=rank_template,
                    parser=parser_used,
                    source=(json_output or {}).get("source"),
                    rankings=template_rankings,
                    best_profile=best_template_profile,
                )
            except Exception:
                template_rankings = []
                best_template_profile = None

        structured_extraction: Optional[Dict[str, Any]] = None
        structured_error: Optional[str] = None
        smart_model_used: Optional[str] = None
        llm_attempts: List[Dict[str, Any]] = []
        llm_ranking: List[Dict[str, Any]] = []
        llm_prompt_compaction: Optional[Dict[str, Any]] = None
        if direct_single_call_used:
            structured_extraction = direct_structured_extraction
            smart_model_used = direct_smart_model_used
            llm_attempts = list(direct_llm_attempts)
            llm_ranking = list(direct_llm_attempts)
            best_llm_score = _best_llm_success_score(llm_ranking or llm_attempts)
            if 0 < best_llm_score < _SMART_ACCEPT_SCORE_THRESHOLD:
                structured_error = (
                    f"low-confidence single-call extraction score {best_llm_score:.2f} "
                    f"(threshold {_SMART_ACCEPT_SCORE_THRESHOLD:.2f})"
                )
                structured_extraction = None
        elif smart_enabled:
            try:
                smart_result = await _run_smart_template_extraction(
                    markdown=extraction_markdown,
                    template_name=resolved_template,
                    model_name=(smart_model or _DOCLING_PRIMARY_MODEL).strip(),
                )
                structured_extraction = smart_result.get("data")
                smart_model_used = smart_result.get("smart_model_used")
                llm_attempts = smart_result.get("llm_attempts") or []
                llm_ranking = smart_result.get("llm_ranking") or []
                llm_prompt_compaction = smart_result.get("prompt_compaction")
                best_llm_score = _best_llm_success_score(llm_ranking or llm_attempts)
                if 0 < best_llm_score < _SMART_ACCEPT_SCORE_THRESHOLD:
                    structured_error = (
                        f"low-confidence LLM extraction score {best_llm_score:.2f} "
                        f"(threshold {_SMART_ACCEPT_SCORE_THRESHOLD:.2f})"
                    )
                    structured_extraction = None
            except Exception as exc:
                structured_error = str(exc)
                if isinstance(exc, SmartExtractionError):
                    llm_attempts = exc.attempts or []
                    llm_ranking = sorted(
                        llm_attempts,
                        key=lambda x: (bool(x.get("success")), float(x.get("score", 0.0))),
                        reverse=True,
                    )
                if resolved_template == "invoice":
                    fallback_data = _extract_invoice_regex_fallback(markdown)
                    if _has_meaningful_invoice_data(fallback_data):
                        structured_extraction = fallback_data
                        smart_model_used = "regex_table_fallback"
                        if structured_error:
                            structured_error = f"{structured_error} | used regex_table_fallback"
                elif resolved_template == "letter_of_credit":
                    fallback_data = _extract_lc_regex_fallback(markdown)
                    if _has_meaningful_lc_data(fallback_data):
                        structured_extraction = fallback_data
                        smart_model_used = "regex_lc_fallback"
                        if structured_error:
                            structured_error = f"{structured_error} | used regex_lc_fallback"
                if smart_required and structured_extraction is None:
                    if smart_enabled:
                        try:
                            _append_llm_ranking_record(
                                template=resolved_template,
                                parser=parser_used,
                                source=(json_output or {}).get("source"),
                                attempts=(llm_seed_ranking or llm_seed_attempts) + (llm_ranking or llm_attempts),
                                selected_model=smart_model_used,
                                status="failed",
                                error=structured_error,
                            )
                            _append_learning_history_record(
                                template=resolved_template,
                                parser=parser_used,
                                smart_model_used=smart_model_used,
                                json_output=json_output,
                                extraction=structured_extraction,
                                status="failed",
                                error=structured_error,
                            )
                        except Exception:
                            pass
                    raise HTTPException(
                        status_code=422,
                        detail=f"Smart extraction failed for template '{resolved_template}': {structured_error}",
                    ) from exc

        if resolved_template == "invoice" and structured_extraction is None:
            fallback_data = _extract_invoice_regex_fallback(markdown)
            if _has_meaningful_invoice_data(fallback_data):
                structured_extraction = fallback_data
                smart_model_used = smart_model_used or "regex_table_fallback"
        elif resolved_template == "letter_of_credit" and structured_extraction is None:
            fallback_data = _extract_lc_regex_fallback(markdown)
            if _has_meaningful_lc_data(fallback_data):
                structured_extraction = fallback_data
                smart_model_used = smart_model_used or "regex_lc_fallback"

        if smart_enabled and structured_extraction is None:
            profile_fallback = _best_template_profile_extraction(resolved_template)
            if profile_fallback and isinstance(profile_fallback.get("extraction"), dict):
                structured_extraction = profile_fallback.get("extraction")
                # Explicitly mark actual resolver used.
                smart_model_used = "template_profile_fallback"
                if not best_template_profile:
                    best_template_profile = {
                        "profile_id": profile_fallback.get("profile_id"),
                        "template": resolved_template,
                        "match_score": 100.0,
                        "match_tier": "99% match",
                    }
                if structured_error:
                    structured_error = f"{structured_error} | used template_profile_fallback"

        if smart_enabled and structured_extraction is not None:
            schema = _schema_for_template_name(resolved_template)
            if schema is not None:
                structured_extraction = _normalize_to_schema(structured_extraction, schema)

        if smart_enabled and structured_extraction:
            learning_template = resolved_template
            try:
                learning_score = _score_learning_record(learning_template, structured_extraction)
                profile_id = _upsert_template_profile(
                    template=learning_template,
                    extraction=structured_extraction,
                    source=(json_output or {}).get("source"),
                    match_score=learning_score,
                )
                if profile_id and best_template_profile and not best_template_profile.get("profile_id"):
                    best_template_profile["profile_id"] = profile_id
                _append_excellent_match_record(
                    template=learning_template,
                    source=(json_output or {}).get("source"),
                    profile_id=profile_id,
                    match_score=learning_score,
                    extraction=structured_extraction,
                )
                _append_learning_raw_sample(
                    template=learning_template,
                    parser=parser_used,
                    smart_model_used=smart_model_used,
                    markdown=markdown,
                    structured_extraction=structured_extraction,
                    json_output=json_output,
                )
                _learn_from_structured_extraction(
                    template=learning_template,
                    markdown=markdown,
                    extraction=structured_extraction,
                    smart_model_used=smart_model_used,
                    json_output=json_output,
                )
            except Exception:
                # Never block conversion due to learning writes.
                pass

        if smart_enabled:
            learning_template = resolved_template
            try:
                status = "success" if structured_extraction else "partial"
                _append_llm_ranking_record(
                    template=learning_template,
                    parser=parser_used,
                    source=(json_output or {}).get("source"),
                    attempts=(llm_seed_ranking or llm_seed_attempts) + (llm_ranking or llm_attempts),
                    selected_model=smart_model_used or llm_seed_model_used,
                    status=status,
                    error=structured_error,
                )
                _append_learning_history_record(
                    template=learning_template,
                    parser=parser_used,
                    smart_model_used=smart_model_used,
                    json_output=json_output,
                    extraction=structured_extraction,
                    status=status,
                    error=structured_error,
                )
            except Exception:
                pass

        primary_model_used = smart_model_used or llm_seed_model_used or used_model
        output_chars = len(markdown)
        if direct_single_call_used and isinstance(structured_extraction, dict):
            output_chars = len(json.dumps(structured_extraction, ensure_ascii=False))
        return {
            "success": True,
            "source": source if source else (file.filename if file else None),
            "extractor": extractor_mode,
            "parser": parser_used,
            "quality": used_quality,
            "fallback_used": fallback_used,
            # Keep "model_used" as the primary pipeline model (smart extraction first, OCR fallback second)
            # so UI messaging is aligned with user-selected LLM when smart extraction is enabled.
            "model_used": primary_model_used,
            "actual_model_used": primary_model_used,
            # Explicit OCR/VL-stage model for transparency in mixed pipelines.
            "ocr_model_used": used_model,
            "markdown": markdown,
            "json_output": json_output,
            "smart_template": resolved_template if smart_enabled else smart_name,
            "smart_model": (smart_model or _DOCLING_PRIMARY_MODEL).strip(),
            "smart_single_call": bool(smart_single_call),
            "seed_stage_skipped": bool(seed_stage_skipped or direct_single_call_used),
            "smart_model_used": smart_model_used,
            "llm_seed_model_used": llm_seed_model_used,
            "llm_seed_prompt_compaction": llm_seed_prompt_compaction,
            "llm_ranking": llm_ranking,
            "llm_prompt_compaction": llm_prompt_compaction or direct_prompt_compaction,
            "template_rankings": template_rankings,
            "best_template_profile": best_template_profile,
            "structured_extraction": structured_extraction,
            "structured_extraction_error": structured_error,
            "single_call_path": bool(direct_single_call_used or (smart_enabled and smart_single_call)),
            "chars": output_chars,
        }
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Docling conversion failed: {exc}") from exc
    finally:
        if file is not None:
            await file.close()
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except OSError:
                pass


@router.post("/refine")
async def refine_markdown(req: RefineRequest) -> Dict[str, Any]:
    instruction = (req.instruction or "").strip()
    markdown = (req.markdown or "").strip()
    if not instruction:
        raise HTTPException(status_code=400, detail="instruction is required")
    if not markdown:
        raise HTTPException(status_code=400, detail="markdown is required")

    chunk_size = max(2000, min(int(req.chunk_size_chars or 10000), 30000))
    chunks = _chunk_text(markdown, max_chars=chunk_size)
    if not chunks:
        raise HTTPException(status_code=400, detail="markdown is empty after normalization")

    engine = get_ai_engine()
    refined_parts: List[str] = []
    template_key = _resolve_template_name(req.template) if req.template else None
    template_schema = _schema_for_template_name(template_key or "") if template_key else None
    template_guidance = _template_guidance(template_key) if template_key and template_schema is not None else ""
    profile_context = ""
    if _INCLUDE_EXEMPLAR_VALUES_IN_PROMPTS and template_key and template_schema is not None:
        try:
            top_profile = _rank_template_profiles(template_key, markdown, top_k=1)
            if top_profile:
                profiles = _load_template_profiles()
                pid = top_profile[0].get("profile_id")
                best = next((p for p in profiles if str(p.get("id")) == str(pid)), None)
                if isinstance(best, dict) and isinstance(best.get("best_extraction"), dict):
                    profile_context = json.dumps(best.get("best_extraction"), ensure_ascii=False, indent=2)
        except Exception:
            profile_context = ""

    for idx, chunk in enumerate(chunks, start=1):
        dynamic_block = ""
        if template_key and template_schema is not None:
            dynamic_block = (
                f"Template: {template_key}\n"
                f"{template_guidance}\n"
                "Target schema:\n"
                f"{json.dumps(template_schema, ensure_ascii=False, indent=2)}\n"
            )
            if profile_context:
                dynamic_block += (
                    "Reference high-match extraction example (use only as structural guidance, "
                    "do not invent missing facts):\n"
                    f"{profile_context}\n"
                )
            else:
                dynamic_block += (
                    "Prompt policy: keys-only schema guidance. "
                    "Do not use exemplar values.\n"
                )
        prompt = (
            "You are refining extracted document markdown.\n"
            "Preserve factual content. Improve clarity, structure, and consistency.\n"
            "Do not invent details.\n\n"
            f"{dynamic_block}\n"
            f"Instruction:\n{instruction}\n\n"
            f"Chunk {idx}/{len(chunks)}:\n{chunk}\n"
        )
        out = await engine.generate_content(req.model, prompt, response_format="text")
        refined_parts.append(str(out))

    if len(refined_parts) == 1:
        final_text = refined_parts[0]
    else:
        merge_prompt = (
            "Merge and normalize the following refined chunk outputs into one cohesive document.\n"
            "Remove duplicates, keep section hierarchy, and keep only grounded facts.\n\n"
            + "\n\n---\n\n".join(refined_parts)
        )
        final_text = str(await engine.generate_content(req.model, merge_prompt, response_format="text"))

    superpowers: Optional[Any] = None
    if req.use_superpowers:
        try:
            from components.superpowers import SuperpowersService  # type: ignore

            sp = SuperpowersService()
            mode = (req.superpower_mode or "review").strip().lower()
            if mode == "plan":
                superpowers = await sp.create_plan(final_text, req.model)
            elif mode == "brainstorm":
                superpowers = await sp.brainstorm(req.instruction, final_text[:12000], req.model)
            elif mode == "debug":
                superpowers = await sp.systematic_debug(final_text, "", req.model)
            else:
                superpowers = await sp.review_code(final_text, req.instruction, req.model)
        except Exception as exc:
            superpowers = {"error": f"superpowers unavailable: {exc}"}

    return {
        "success": True,
        "chunks": len(chunks),
        "refined_markdown": final_text,
        "superpowers": superpowers,
    }


@router.get("/templates/definitions")
async def list_template_definitions() -> Dict[str, Any]:
    defs = _load_template_definitions()
    rows: List[Dict[str, Any]] = []
    for name, cfg in defs.items():
        schema = cfg.get("schema") or {}
        rows.append(
            {
                "name": name,
                "aliases": cfg.get("aliases") or [],
                "guidance": cfg.get("guidance") or "",
                "field_count": len(schema.keys()) if isinstance(schema, dict) else 0,
                "schema": schema,
            }
        )
    rows = sorted(rows, key=lambda x: x["name"])
    return {"success": True, "templates": rows, "count": len(rows)}


@router.post("/templates/definitions")
async def upsert_template_definition(req: TemplateDefinitionRequest) -> Dict[str, Any]:
    name = _normalize_template_key(req.name)
    if not name or name in {"all", "none"}:
        raise HTTPException(status_code=400, detail="name must be a valid template identifier")
    if not isinstance(req.schema_payload, dict) or not req.schema_payload:
        raise HTTPException(status_code=400, detail="schema must be a non-empty JSON object")

    defs = _load_template_definitions()
    aliases: List[str] = [name]
    for raw in (req.aliases or []):
        a = _normalize_template_key(raw)
        if a and a not in aliases:
            aliases.append(a)
    defs[name] = {
        "schema": req.schema_payload,
        "guidance": (req.guidance or "").strip(),
        "aliases": aliases,
    }
    _save_template_definitions(defs)
    _append_template_upload_history(
        action="upsert_definition",
        template=name,
        source="api",
        details={"aliases": aliases, "field_count": len(req.schema_payload.keys())},
    )
    return {"success": True, "template": name, "aliases": aliases}


@router.delete("/templates/definitions/{template_name}")
async def delete_template_definition(template_name: str, delete_profiles: bool = True) -> Dict[str, Any]:
    name = _normalize_learning_template(template_name)
    if not name or name == "all":
        raise HTTPException(status_code=400, detail="Invalid template name")
    defs = _load_template_definitions()
    if name not in defs:
        raise HTTPException(status_code=404, detail=f"Template '{template_name}' not found")
    defs.pop(name, None)
    _save_template_definitions(defs)

    removed_profiles = 0
    if delete_profiles:
        profiles = _load_template_profiles()
        kept = []
        for p in profiles:
            t = _normalize_learning_template(p.get("template"))
            if t == name:
                removed_profiles += 1
                continue
            kept.append(p)
        _save_template_profiles(kept)

    _append_template_upload_history(
        action="delete_definition",
        template=name,
        source="api",
        details={"delete_profiles": bool(delete_profiles), "removed_profiles": removed_profiles},
    )
    return {"success": True, "template": name, "removed_profiles": removed_profiles}


@router.post("/templates/exemplars")
async def upsert_template_exemplar(req: TemplateExemplarRequest) -> Dict[str, Any]:
    exemplar = req.exemplar if isinstance(req.exemplar, dict) else {}
    if not exemplar:
        raise HTTPException(status_code=400, detail="exemplar must be a non-empty JSON object")
    template_name = _normalize_template_key(req.template) if req.template else _guess_template_name_from_payload(exemplar)
    if not template_name:
        raise HTTPException(status_code=400, detail="template name is required")

    canonical = _register_template_from_example(
        template_name=template_name,
        exemplar=exemplar,
        guidance=(req.guidance or "").strip(),
        aliases=req.aliases or [],
    )
    schema = _schema_for_template_name(canonical)
    normalized_extraction = _normalize_to_schema(exemplar, schema) if schema is not None else exemplar
    score = max(_EXCELLENT_MATCH_THRESHOLD, _score_learning_record(canonical, normalized_extraction))
    profile_id = _upsert_template_profile(
        template=canonical,
        extraction=normalized_extraction,
        source=(req.source or "manual").strip() or "manual",
        match_score=score,
    )
    _append_excellent_match_record(
        template=canonical,
        source=(req.source or "manual").strip() or "manual",
        profile_id=profile_id,
        match_score=score,
        extraction=normalized_extraction,
    )
    _append_template_upload_history(
        action="upsert_exemplar",
        template=canonical,
        source=(req.source or "manual").strip() or "manual",
        details={"profile_id": profile_id, "score": round(float(score), 2)},
    )
    return {
        "success": True,
        "template": canonical,
        "profile_id": profile_id,
        "score": round(float(score), 2),
        "permanent": True,
    }


@router.post("/templates/upload")
async def upload_template_file(
    file: UploadFile = File(...),
    template: Optional[str] = None,
    guidance: Optional[str] = None,
) -> Dict[str, Any]:
    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Empty file")
    try:
        data = json.loads(raw.decode("utf-8"))
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Invalid JSON file: {exc}") from exc

    name = _normalize_template_key(template) if template else ""
    if isinstance(data, dict) and isinstance(data.get("schema"), dict):
        def_name = name or _normalize_template_key(data.get("name")) or "custom_template"
        aliases = data.get("aliases") if isinstance(data.get("aliases"), list) else []
        defs = _load_template_definitions()
        defs[def_name] = {
            "schema": data.get("schema"),
            "guidance": (guidance or data.get("guidance") or "").strip(),
            "aliases": [def_name] + [a for a in [_normalize_template_key(x) for x in aliases] if a and a != def_name],
        }
        _save_template_definitions(defs)
        _append_template_upload_history(
            action="upload_definition_file",
            template=def_name,
            source=(file.filename or "upload.json"),
            details={"field_count": len((data.get("schema") or {}).keys())},
        )
        return {"success": True, "mode": "definition", "template": def_name}

    exemplar = data if isinstance(data, dict) else {}
    req = TemplateExemplarRequest(
        template=(name or None),
        exemplar=exemplar,
        guidance=(guidance or ""),
        source=(file.filename or "upload.json"),
    )
    result = await upsert_template_exemplar(req)
    _append_template_upload_history(
        action="upload_exemplar_file",
        template=result.get("template") or (name or "unknown"),
        source=(file.filename or "upload.json"),
        details={"profile_id": result.get("profile_id"), "score": result.get("score")},
    )
    return {"success": True, "mode": "exemplar", **result}


@router.post("/learning/reset")
async def reset_learning(template: str = "all", clear_raw: bool = True) -> Dict[str, Any]:
    normalized_template = _normalize_learning_template(template)
    result = _reset_learning_state(template=normalized_template, clear_raw=clear_raw)
    return {
        "success": True,
        **result,
    }


@router.get("/learning/history")
async def learning_history(limit: int = 200, template: str = "all") -> Dict[str, Any]:
    return {
        "success": True,
        "template": _normalize_learning_template(template),
        **_load_learning_history(limit=limit, template=template),
    }


@router.get("/learning/excellent-matches")
async def excellent_matches_history(limit: int = 200, template: str = "all") -> Dict[str, Any]:
    path = _excellent_match_history_path()
    normalized_template = _normalize_learning_template(template)
    if not os.path.exists(path):
        return {"success": True, "template": normalized_template, "records": [], "count": 0}
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        lines = [ln.strip() for ln in f if ln.strip()]
    rows: List[Dict[str, Any]] = []
    for ln in reversed(lines):
        try:
            rec = json.loads(ln)
        except Exception:
            continue
        t = _normalize_learning_template(rec.get("template"))
        if normalized_template != "all" and t != normalized_template:
            continue
        rows.append(
            {
                "ts": rec.get("ts"),
                "template": t,
                "source": rec.get("source"),
                "profile_id": rec.get("profile_id"),
                "match_score": rec.get("match_score"),
                "extraction": rec.get("extraction"),
            }
        )
        if len(rows) >= max(1, min(int(limit), 1000)):
            break
    return {"success": True, "template": normalized_template, "records": rows, "count": len(rows)}


@router.get("/templates/uploads/history")
async def template_uploads_history(limit: int = 200, template: str = "all") -> Dict[str, Any]:
    path = _template_upload_history_path()
    normalized_template = _normalize_learning_template(template)
    if not os.path.exists(path):
        return {"success": True, "template": normalized_template, "records": [], "count": 0}
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        lines = [ln.strip() for ln in f if ln.strip()]
    rows: List[Dict[str, Any]] = []
    for ln in reversed(lines):
        try:
            rec = json.loads(ln)
        except Exception:
            continue
        t = _normalize_learning_template(rec.get("template"))
        if normalized_template != "all" and t != normalized_template:
            continue
        rows.append(
            {
                "ts": rec.get("ts"),
                "action": rec.get("action"),
                "template": t,
                "source": rec.get("source"),
                "details": rec.get("details") or {},
            }
        )
        if len(rows) >= max(1, min(int(limit), 1000)):
            break
    return {"success": True, "template": normalized_template, "records": rows, "count": len(rows)}


def _load_template_ranking_history(limit: int = 200, template: str = "all") -> Dict[str, Any]:
    path = _template_ranking_history_path()
    normalized_template = _normalize_learning_template(template)
    if not os.path.exists(path):
        return {"records": [], "summary": {"total": 0, "avg_best_match": 0.0, "by_tier": {"99% match": 0, "95% match": 0, "90% match": 0, "<90% match": 0}}}

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        lines = [ln.strip() for ln in f if ln.strip()]

    records: List[Dict[str, Any]] = []
    for ln in reversed(lines):
        try:
            rec = json.loads(ln)
        except Exception:
            continue
        t = _normalize_learning_template(rec.get("template"))
        if normalized_template != "all" and t != normalized_template:
            continue
        best = rec.get("best_profile") or {}
        score = float(best.get("match_score") or 0.0)
        records.append(
            {
                "ts": rec.get("ts"),
                "template": t,
                "parser": rec.get("parser"),
                "source": rec.get("source"),
                "best_profile_id": best.get("profile_id"),
                "match_score": score,
                "match_tier": _tier_from_score(score),
                "rankings": rec.get("rankings") or [],
            }
        )
        if len(records) >= max(1, min(int(limit), 1000)):
            break

    tiers = {"99% match": 0, "95% match": 0, "90% match": 0, "<90% match": 0}
    for r in records:
        tiers[r["match_tier"]] = tiers.get(r["match_tier"], 0) + 1
    avg_best = round(sum(float(r["match_score"]) for r in records) / len(records), 2) if records else 0.0
    return {"records": records, "summary": {"total": len(records), "avg_best_match": avg_best, "by_tier": tiers}}


@router.get("/learning/template-rankings/history")
async def template_rankings_history(limit: int = 200, template: str = "all") -> Dict[str, Any]:
    return {
        "success": True,
        "template": _normalize_learning_template(template),
        **_load_template_ranking_history(limit=limit, template=template),
    }


@router.get("/learning/template-profiles")
async def template_profiles(template: str = "all", limit: int = 200) -> Dict[str, Any]:
    normalized_template = _normalize_learning_template(template)
    profiles = _load_template_profiles()
    rows: List[Dict[str, Any]] = []
    for p in profiles:
        t = _normalize_learning_template(p.get("template"))
        if normalized_template != "all" and t != normalized_template:
            continue
        rows.append(
            {
                "id": p.get("id"),
                "template": t,
                "seen_count": p.get("seen_count"),
                "permanent": bool(p.get("permanent")),
                "avg_match_score": p.get("avg_match_score"),
                "avg_json_score": p.get("avg_json_score"),
                "best_json_score": p.get("best_json_score"),
                "last_json_score": p.get("last_json_score"),
                "sample_source": p.get("sample_source"),
                "updated_at": p.get("updated_at"),
            }
        )
    rows = sorted(
        rows,
        key=lambda x: (
            float(x.get("best_json_score") or x.get("avg_json_score") or x.get("avg_match_score") or 0.0),
            float(x.get("avg_json_score") or x.get("avg_match_score") or 0.0),
            int(x.get("seen_count") or 0),
        ),
        reverse=True,
    )
    rows = rows[: max(1, min(int(limit), 1000))]
    return {"success": True, "template": normalized_template, "profiles": rows, "count": len(rows)}


@router.get("/learning/llm-ranking/history")
async def llm_ranking_history(limit: int = 200, template: str = "all") -> Dict[str, Any]:
    model_metrics: Dict[str, Any] = {}
    try:
        engine = get_ai_engine()
        metrics = await engine.get_model_call_metrics(days=7, status="success")
        model_metrics = metrics.get("models", {}) if isinstance(metrics, dict) else {}
    except Exception:
        model_metrics = {}
    return {
        "success": True,
        "template": _normalize_learning_template(template),
        "model_metrics": model_metrics,
        **_load_llm_ranking_history(limit=limit, template=template),
    }


@router.get("/learning/model-health")
async def learning_model_health() -> Dict[str, Any]:
    data = _load_model_health()
    return {
        "success": True,
        "phase_prefs": data.get("phase_prefs", {}),
        "stats": data.get("stats", {}),
        "rankings": data.get("rankings", {}),
    }


def setup(mcp=None, app=None):
    if app is not None:
        app.include_router(router)

    if mcp is not None:

        @mcp.tool()
        async def docling_convert_to_markdown(source: str, quality: str = "fast") -> str:
            """Convert a local file path or URL to markdown using Docling."""
            return _convert_source_to_markdown(source, quality=quality)
